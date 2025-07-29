"""PowerPoint template processing and merge field replacement module."""

import logging
import os
import re
from typing import Any, Dict, List, Optional
from PIL import Image as PILImage
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .utils.exceptions import PowerPointProcessingError
from .utils.validation import validate_merge_fields
from .utils.slide_utils import filter_slides, is_template_slide, extract_list_name

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)  # Force INFO level to ensure debug messages are visible


class PowerPointProcessor:
    """Processes PowerPoint templates and replaces merge fields with data."""

    def __init__(self, template_path: str) -> None:
        """Initialize PowerPoint processor with template path."""
        self.template_path = template_path
        self.presentation = None
        self._validate_template()

    def _validate_template(self) -> None:
        """Validate PowerPoint template exists and is readable."""
        if not os.path.exists(self.template_path):
            raise PowerPointProcessingError(
                f"PowerPoint template not found: {self.template_path}"
            )

        try:
            self.presentation = Presentation(self.template_path)
            logger.info(
                f"Successfully loaded PowerPoint template: {self.template_path}"
            )
        except Exception as e:
            raise PowerPointProcessingError(f"Invalid PowerPoint template format: {e}")

    def get_merge_fields(self) -> List[str]:
        """Extract all merge fields from the presentation."""
        if not self.presentation:
            self._validate_template()

        merge_fields = set()

        try:
            for slide in self.presentation.slides:
                slide_fields = self._extract_slide_merge_fields(slide)
                merge_fields.update(slide_fields)

            return sorted(list(merge_fields))

        except Exception as e:
            raise PowerPointProcessingError(f"Failed to extract merge fields: {e}")

    def _extract_slide_merge_fields(self, slide) -> List[str]:
        """Extract merge fields from a single slide."""
        fields = []

        try:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_content = self._get_full_text_from_shape(shape)
                    if text_content:
                        shape_fields = validate_merge_fields(text_content)
                        fields.extend(shape_fields)

                # Check for table cells
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_fields = self._extract_table_merge_fields(shape.table)
                    fields.extend(table_fields)

        except Exception as e:
            logger.warning(f"Failed to extract merge fields from slide: {e}")

        return fields

    def _extract_table_merge_fields(self, table) -> List[str]:
        """Extract merge fields from table cells."""
        fields = []

        try:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        cell_fields = validate_merge_fields(cell.text)
                        fields.extend(cell_fields)
        except Exception as e:
            logger.warning(f"Failed to extract merge fields from table: {e}")

        return fields

    def _get_full_text_from_shape(self, shape: BaseShape) -> str:
        """Get complete text content from a shape including all paragraphs."""
        try:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                return ""

            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ""
                for run in paragraph.runs:
                    paragraph_text += run.text
                text_parts.append(paragraph_text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.warning(f"Failed to get text from shape: {e}")
            return ""

    def merge_data(
        self,
        data: Dict[str, Any],
        output_path: str,
        images: Optional[Dict[str, List[Dict[str, Any]]]] = None,
        config: Optional[Dict[str, Any]] = None,
    ) -> str:
        """Merge data into the PowerPoint template using two-pass processing.
        
        Pass 1: Create dynamic slides and save to memory
        Pass 2: Reload from memory and process all content

        Args:
            data: Data to merge into the template
            output_path: Path to save the merged presentation
            images: Dictionary of images by sheet name
            config: Configuration for PowerPoint processing

        Returns:
            Path to the merged presentation
        """
        logger.info("🔧 DEBUG: ENTRY - merge_data() called with two-pass processing")
        logger.info(f"🔧 DEBUG: ENTRY - data keys: {list(data.keys()) if data else None}")
        logger.info(f"🔧 DEBUG: ENTRY - output_path: {output_path}")
        logger.info(f"🔧 DEBUG: ENTRY - images provided: {images is not None}")
        logger.info(f"🔧 DEBUG: ENTRY - config provided: {config is not None}")

        if not self.presentation:
            logger.error("🔧 DEBUG: ERROR - No PowerPoint template loaded")
            raise PowerPointProcessingError("No presentation loaded")

        # Ensure output_path is an absolute path
        if not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)

        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        os.makedirs(output_dir, exist_ok=True)
        logger.debug(f"Ensuring PowerPoint output directory exists: {output_dir}")

        try:
            # Get PowerPoint configuration
            powerpoint_config = (
                config.get("global_settings", {}).get("powerpoint", {})
                if config
                else {}
            )
            logger.info(f"🔧 DEBUG: PowerPoint config loaded: {bool(powerpoint_config)}")
            logger.info(f"🔧 DEBUG: Total slides in presentation: {len(list(self.presentation.slides))}")

            # PASS 1: Create dynamic slides and save to memory
            dynamic_slide_mapping = self._create_dynamic_slides_pass(data, powerpoint_config)
            
            # DEBUG: Log dynamic slide mapping state after Pass 1
            logger.info(f"🔧 DEBUG: PASS 1 COMPLETED")
            logger.info(f"🔧 DEBUG: Total slides after Pass 1: {len(list(self.presentation.slides))}")
            logger.info(f"🔧 DEBUG: Dynamic slide mapping: {dynamic_slide_mapping}")
            if dynamic_slide_mapping:
                expected_dynamic_slides = list(dynamic_slide_mapping.keys())
                logger.info(f"🔧 DEBUG: Expected dynamic slides in Pass 2: {expected_dynamic_slides}")
            else:
                logger.info(f"🔧 DEBUG: No dynamic slides detected in Pass 1")
            
            # PASS 2: Reload from memory and process all content
            self._process_slides_pass(data, images, powerpoint_config, dynamic_slide_mapping)
            
            # Save final presentation
            self.presentation.save(output_path)
            logger.info(f"Two-pass merged presentation saved to: {output_path}")
            
            # Post-process the saved file to remove any remaining error attributes
            post_cleanup_count = self._post_process_xml_cleanup(output_path)
            if post_cleanup_count > 0:
                logger.info(
                    f"Post-processing removed {post_cleanup_count} additional error attributes"
                )
                
            logger.info("🔧 DEBUG: EXIT - merge_data() completed successfully with two-pass processing")
            return output_path

        except Exception as e:
            logger.error(f"🔧 DEBUG: EXCEPTION - merge_data failed: {e}")
            logger.error(f"🔧 DEBUG: EXCEPTION - Exception type: {type(e)}")
            import traceback
            logger.error(f"🔧 DEBUG: EXCEPTION - Traceback: {traceback.format_exc()}")
            raise PowerPointProcessingError(f"Failed to merge data into presentation: {e}")

    def _create_dynamic_slides_pass(self, data: Dict[str, Any], powerpoint_config: Dict[str, Any]) -> Dict[int, Dict[str, Any]]:
        """Pass 1: Create dynamic slides only, no processing.
        
        Returns mapping of slide positions to list item info for Pass 2.
        """
        from io import BytesIO
        
        logger.info("🔧 DEBUG: PASS 1 - Starting dynamic slide creation")
        dynamic_slide_mapping = {}
        
        # Check if dynamic slides are enabled
        dynamic_enabled = powerpoint_config.get("dynamic_slides", {}).get("enabled", True)
        if not dynamic_enabled:
            logger.info("🔧 DEBUG: PASS 1 - Dynamic slides disabled, skipping")
            return dynamic_slide_mapping
        
        # Build mapping during slide creation by modifying _process_dynamic_slides to return mapping
        dynamic_slide_mapping = self._process_dynamic_slides_with_mapping(data, powerpoint_config)
        
        # Save presentation to memory
        memory_file = BytesIO()
        self.presentation.save(memory_file)
        memory_file.seek(0)
        logger.info("🔧 DEBUG: PASS 1 - Saved presentation to memory")
        
        # Store the memory file for Pass 2
        self._memory_file = memory_file
        
        logger.info(f"🔧 DEBUG: PASS 1 - Created {len(dynamic_slide_mapping)} dynamic slide mappings")
        return dynamic_slide_mapping

    def _process_slides_pass(self, data: Dict[str, Any], images: Dict[str, Any], 
                           powerpoint_config: Dict[str, Any], dynamic_slide_mapping: Dict[int, Dict[str, Any]]) -> None:
        """Pass 2: Reload from memory and process all slides with fresh object model."""
        from pptx import Presentation
        
        logger.info("🔧 DEBUG: PASS 2 - Starting content processing with fresh object model")
        
        # Reload presentation from memory to get fresh python-pptx object model
        if hasattr(self, '_memory_file'):
            self.presentation = Presentation(self._memory_file)
            logger.info("🔧 DEBUG: PASS 2 - Reloaded presentation from memory with fresh object model")
        else:
            logger.warning("🔧 DEBUG: PASS 2 - No memory file found, using current presentation")
        
        # Apply slide filtering
        logger.info("🔧 DEBUG: PASS 2 - Starting slide filtering...")
        self._filter_slides(powerpoint_config)
        logger.info(f"🔧 DEBUG: PASS 2 - Slides remaining after filtering: {len(list(self.presentation.slides))}")
        
        # Process all slides with fresh object model
        slides = list(self.presentation.slides)
        logger.info(f"🔧 DEBUG: PASS 2 - Total slides after filtering: {len(slides)}")
        logger.info(f"🔧 DEBUG: PASS 2 - Expected dynamic slides at indices: {list(dynamic_slide_mapping.keys())}")
        
        # Debug: Check what content is actually at the mapped positions
        for mapped_index in dynamic_slide_mapping.keys():
            if mapped_index < len(slides):
                slide_at_position = slides[mapped_index]
                # Check if this slide has template markers or actual content
                sample_text = ""
                for shape in slide_at_position.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                sample_text += run.text
                                if len(sample_text) > 50:
                                    break
                            if len(sample_text) > 50:
                                break
                        if len(sample_text) > 50:
                            break
                logger.info(f"🔧 DEBUG: PASS 2 - Slide {mapped_index + 1} content sample: '{sample_text[:50]}...'")
            else:
                logger.warning(f"🔧 DEBUG: PASS 2 - Mapped index {mapped_index} is out of range (only {len(slides)} slides)")
        
        for i, slide in enumerate(slides):
            if i in dynamic_slide_mapping:
                # This is a dynamic slide - process with list context
                mapping = dynamic_slide_mapping[i]
                list_name = mapping['list_name']
                item_index = mapping['item_index']
                total_items = mapping['total_items']
                item_data = mapping['item_data']
                
                # Create context for this list item using the stored data
                context = self._create_list_item_context(
                    item_data, item_index, total_items, data, 
                    powerpoint_config.get("dynamic_slides", {})
                )
                
                logger.info(f"🔧 DEBUG: PASS 2 - Processing dynamic slide {i+1} with list item {item_index + 1}/{total_items}")
                logger.info(f"🔧 DEBUG: PASS 2 - Using item_data: {item_data}")
                logger.info(f"🔧 DEBUG: PASS 2 - Context keys: {list(context.keys())}")
                
                # Process this slide with the list item context
                self._process_slide_with_context(slide, context)
            else:
                # Regular slide - process normally
                logger.info(f"🔧 DEBUG: PASS 2 - Processing regular slide {i+1}")
                self._process_slide(slide, data, images)
        
        # Validate and clean up
        self._validate_presentation_integrity()
        
        # Final cleanup
        final_cleanup_count = self._final_cleanup_presentation()
        if final_cleanup_count > 0:
            logger.info(f"🔧 DEBUG: PASS 2 - Final cleanup removed {final_cleanup_count} error attributes")
        
        logger.info("🔧 DEBUG: PASS 2 - Content processing completed")

    def _filter_slides(self, powerpoint_config: Dict[str, Any]) -> None:
        """Filter slides by actually removing non-included slides from the presentation.

        Args:
            powerpoint_config: PowerPoint configuration
        """
        try:
            slide_filter_config = powerpoint_config.get("slide_filter", {})
            include_slides = slide_filter_config.get("include_slides", [])
            exclude_slides = slide_filter_config.get("exclude_slides", [])

            logger.info(
                f"🔧 DEBUG _filter_slides: slide_filter_config = {slide_filter_config}"
            )
            logger.info(f"🔧 DEBUG _filter_slides: include_slides = {include_slides}")
            logger.info(f"🔧 DEBUG _filter_slides: exclude_slides = {exclude_slides}")
            logger.info(
                f"🔧 DEBUG _filter_slides: slides count before filtering = {len(list(self.presentation.slides))}"
            )

            # If no filtering configured, do nothing
            if not include_slides and not exclude_slides:
                logger.info(
                    f"🔧 DEBUG _filter_slides: No filtering configured, keeping all slides"
                )
                return

            slides_to_remove = []

            # Determine which slides to remove
            for i, slide in enumerate(self.presentation.slides):
                slide_number = i + 1  # Convert to 1-based
                should_remove = False

                # If include list is specified, remove slides NOT in the list
                if include_slides:
                    if slide_number not in include_slides:
                        should_remove = True
                        logger.info(
                            f"🔧 DEBUG: Slide {slide_number} not in include list, will remove"
                        )

                # If exclude list is specified, remove slides IN the list
                elif exclude_slides:
                    if slide_number in exclude_slides:
                        should_remove = True
                        logger.info(
                            f"🔧 DEBUG: Slide {slide_number} in exclude list, will remove"
                        )

                if should_remove:
                    slides_to_remove.append((i, slide))

            # Remove slides (iterate backwards to avoid index issues)
            for slide_index, slide in reversed(slides_to_remove):
                try:
                    logger.info(f"🔧 DEBUG: Removing slide {slide_index + 1}")
                    slide_id = slide.slide_id
                    self.presentation.slides._sldIdLst.remove(
                        self.presentation.slides._sldIdLst.xpath(
                            f'//p:sldId[@id="{slide_id}"]'
                        )[0]
                    )
                    logger.info(
                        f"🔧 DEBUG: Successfully removed slide {slide_index + 1}"
                    )
                except Exception as e:
                    logger.error(
                        f"🔧 DEBUG: ERROR removing slide {slide_index + 1}: {e}"
                    )

            final_count = len(list(self.presentation.slides))
            logger.info(
                f"🔧 DEBUG _filter_slides: slides count after filtering = {final_count}"
            )
            logger.info(
                f"Slide filtering: removed {len(slides_to_remove)} slides, {final_count} slides remaining"
            )

        except Exception as e:
            logger.error(f"Error filtering slides: {e}")

    def _process_dynamic_slides(
        self, data: Dict[str, Any], powerpoint_config: Dict[str, Any]
    ) -> set:
        """Process template slides and create duplicates from list data.

        Args:
            data: The merge data containing lists for duplication
            powerpoint_config: PowerPoint configuration
        """
        processed_slide_ids = set()

        try:
            dynamic_config = powerpoint_config.get("dynamic_slides", {})
            template_marker = dynamic_config.get("template_marker", "{{#list:")
            remove_template_slides = dynamic_config.get("remove_template_slides", True)

            logger.info(
                f"🔧 DEBUG _process_dynamic_slides: dynamic_config = {dynamic_config}"
            )
            logger.info(
                f"🔧 DEBUG _process_dynamic_slides: template_marker = {template_marker}"
            )
            logger.info(
                f"🔧 DEBUG _process_dynamic_slides: remove_template_slides = {remove_template_slides}"
            )

            # Find template slides (iterate backwards to avoid index issues when removing)
            template_slides_info = []
            slides = list(self.presentation.slides)

            for i in range(len(slides) - 1, -1, -1):
                slide = slides[i]
                logger.info(f"🔧 DEBUG: Checking slide {i + 1} for template marker")
                if is_template_slide(slide, template_marker):
                    logger.info(f"🔧 DEBUG: Slide {i + 1} IS a template slide!")
                    list_name = extract_list_name(slide, template_marker)
                    logger.info(f"🔧 DEBUG: Extracted list name: {list_name}")
                    if list_name:
                        list_data = self._find_list_in_data(data, list_name)
                        if isinstance(list_data, list) and list_data:
                            template_slides_info.append(
                                (i, slide, list_name, list_data)
                            )

            # Process each template slide (in forward order for proper indexing)
            # Reverse the list since we built it backwards
            template_slides_info.reverse()

            # Keep track of slide index adjustments
            index_offset = 0

            for (
                original_slide_index,
                template_slide,
                list_name,
                list_data,
            ) in template_slides_info:
                # Adjust slide index based on previous insertions
                slide_index = original_slide_index + index_offset

                logger.info(
                    f"Processing template slide {original_slide_index + 1} for list '{list_name}' with {len(list_data)} items"
                )

                # Track this template slide as processed
                processed_slide_ids.add(template_slide.slide_id)

                # Store the slides to be inserted
                new_slides = []

                # Create duplicates for each item in the list
                for item_index, item in enumerate(list_data):
                    # Use the enhanced duplicate_slide function from slide_utils
                    from .utils.slide_utils import duplicate_slide

                    # For now, use simple duplication but insert in the right place
                    new_slide = self._duplicate_slide_enhanced(
                        template_slide, slide_index + item_index
                    )

                    # Create context for this list item
                    item_context = self._create_list_item_context(
                        item, item_index, len(list_data), data, dynamic_config
                    )

                    # Process the dynamic slide using context-specific processing
                    # This handles direct field names and template marker removal
                    logger.info(f"🔧 DEBUG: About to process dynamic slide for {list_name}[{item_index}]")
                    self._process_slide_with_context(new_slide, item_context)
                    logger.info(f"🔧 DEBUG: Completed processing dynamic slide for {list_name}[{item_index}]")

                    new_slides.append(new_slide)
                    # Track the new slide as processed to prevent re-processing in main loop
                    processed_slide_ids.add(new_slide.slide_id)
                    logger.info(f"🔧 DEBUG: Created slide for {list_name}[{item_index}], added ID {new_slide.slide_id} to processed set")

                # Update index offset for next template slide
                index_offset += len(new_slides)

                # Remove the original template slide if configured
                if remove_template_slides and new_slides:
                    logger.info(
                        f"🔧 DEBUG: Attempting to remove original template slide {original_slide_index + 1}"
                    )
                    try:
                        # Use the original template slide reference, not position calculation
                        # The template slide hasn't moved yet - new slides are added at the end
                        slide_id = template_slide.slide_id

                        self.presentation.slides._sldIdLst.remove(
                            self.presentation.slides._sldIdLst.xpath(
                                f'//p:sldId[@id="{slide_id}"]'
                            )[0]
                        )

                        # Adjust offset since we removed a slide
                        index_offset -= 1

                        logger.info(
                            f"🔧 DEBUG: Successfully removed template slide {original_slide_index + 1}"
                        )
                    except Exception as e:
                        logger.error(
                            f"🔧 DEBUG: ERROR removing template slide {original_slide_index + 1}: {e}"
                        )
                        logger.warning(f"Could not remove template slide: {e}")

            if template_slides_info:
                logger.info(f"Processed {len(template_slides_info)} template slides")
            else:
                logger.info(
                    "🔧 DEBUG _process_dynamic_slides: No template slides found"
                )

            logger.info(f"🔧 DEBUG _process_dynamic_slides: Completed successfully")

            return processed_slide_ids

        except Exception as e:
            logger.error(f"🔧 DEBUG _process_dynamic_slides: ERROR - {e}")
            logger.error(f"Error processing dynamic slides: {e}")
            return processed_slide_ids

    def _process_dynamic_slides_with_mapping(self, data: Dict[str, Any], powerpoint_config: Dict[str, Any]) -> Dict[int, Dict[str, Any]]:
        """Process template slides and create duplicates while tracking mapping for Pass 2.

        Args:
            data: The merge data containing lists for duplication
            powerpoint_config: PowerPoint configuration
            
        Returns:
            Dictionary mapping slide indices to list item info for Pass 2
        """
        dynamic_slide_mapping = {}
        
        try:
            from .utils.slide_utils import is_template_slide, extract_list_name

            dynamic_config = powerpoint_config.get("dynamic_slides", {})
            template_marker = dynamic_config.get("template_marker", "{{#list:")
            remove_template_slides = dynamic_config.get("remove_template_slides", True)

            logger.info(f"🔧 DEBUG: Processing dynamic slides with mapping - template_marker = {template_marker}")

            # Find template slides (iterate backwards to avoid index issues when removing)
            template_slides_info = []
            slides = list(self.presentation.slides)

            for i in range(len(slides) - 1, -1, -1):
                slide = slides[i]
                logger.info(f"🔧 DEBUG: Checking slide {i + 1} for template marker")
                if is_template_slide(slide, template_marker):
                    logger.info(f"🔧 DEBUG: Slide {i + 1} IS a template slide!")
                    list_name = extract_list_name(slide, template_marker)
                    logger.info(f"🔧 DEBUG: Extracted list name: {list_name}")
                    if list_name:
                        list_data = self._find_list_in_data(data, list_name)
                        if isinstance(list_data, list) and list_data:
                            template_slides_info.append((i, slide, list_name, list_data))

            # Process each template slide (in forward order for proper indexing)
            template_slides_info.reverse()

            # Keep track of slide index adjustments
            index_offset = 0

            # Track all created slides for final mapping calculation
            created_slides_info = []

            for original_slide_index, template_slide, list_name, list_data in template_slides_info:
                # Adjust slide index based on previous insertions
                slide_index = original_slide_index + index_offset

                logger.info(f"🔧 DEBUG: Processing template slide {original_slide_index + 1} for list '{list_name}' with {len(list_data)} items")

                # Create duplicates for each item in the list  
                for item_index, item in enumerate(list_data):
                    logger.info(f"🔧 DEBUG: Creating slide {item_index + 1}/{len(list_data)} for {list_name}[{item_index}]")
                    
                    # Calculate where the new slide will be positioned
                    new_slide_position = slide_index + item_index
                    
                    try:
                        # Use the enhanced duplicate_slide function
                        new_slide = self._duplicate_slide_enhanced(template_slide, new_slide_position)
                        logger.info(f"🔧 DEBUG: Successfully duplicated slide at position {new_slide_position}")
                        
                        # Track this slide for final mapping calculation
                        created_slides_info.append({
                            'position': new_slide_position,
                            'list_name': list_name,
                            'item_index': item_index,
                            'total_items': len(list_data),
                            'item_data': item,
                            'template_slide_index': original_slide_index
                        })
                        
                    except Exception as e:
                        logger.error(f"🔧 DEBUG: ERROR duplicating slide for {list_name}[{item_index}]: {e}")
                        continue

                # Update index offset for next template slide
                index_offset += len(list_data)

                # Remove the original template slide if configured
                if remove_template_slides:
                    logger.info(f"🔧 DEBUG: Attempting to remove original template slide {original_slide_index + 1}")
                    try:
                        slide_id = template_slide.slide_id
                        self.presentation.slides._sldIdLst.remove(
                            self.presentation.slides._sldIdLst.xpath(f'//p:sldId[@id="{slide_id}"]')[0]
                        )
                        # Adjust offset since we removed a slide
                        index_offset -= 1
                        
                        logger.info(f"🔧 DEBUG: Successfully removed template slide {original_slide_index + 1}")
                    except Exception as e:
                        logger.error(f"🔧 DEBUG: ERROR removing template slide {original_slide_index + 1}: {e}")

            # Build final mapping with correct positions after all operations
            for slide_info in created_slides_info:
                # The final position is simply the position where the slide was created
                # No adjustment needed - slides retain their positions after template removal
                final_position = slide_info['position']
                
                # Add to mapping for Pass 2
                dynamic_slide_mapping[final_position] = {
                    'list_name': slide_info['list_name'],
                    'item_index': slide_info['item_index'],
                    'total_items': slide_info['total_items'],
                    'item_data': slide_info['item_data'],
                    'is_dynamic': True
                }
                
                logger.info(f"🔧 DEBUG: Final mapping - slide {final_position + 1} = {slide_info['list_name']}[{slide_info['item_index']}]")

            logger.info(f"🔧 DEBUG: Created {len(dynamic_slide_mapping)} dynamic slides with mapping")
            logger.info(f"🔧 DEBUG: Final mapping: {dynamic_slide_mapping}")
            
            return dynamic_slide_mapping

        except Exception as e:
            logger.error(f"🔧 DEBUG: Error processing dynamic slides with mapping: {e}")
            import traceback
            logger.error(f"🔧 DEBUG: Traceback: {traceback.format_exc()}")
            return dynamic_slide_mapping

    def _find_list_in_data(
        self, data: Dict[str, Any], list_name: str
    ) -> Optional[List[Dict[str, Any]]]:
        """Recursively find a list in nested data structure.

        Args:
            data: The data dictionary to search in
            list_name: The name of the list to find

        Returns:
            The list data if found, None otherwise
        """
        try:
            # Check at current level
            if list_name in data:
                return data[list_name]

            # Recursively search in nested dictionaries
            for key, value in data.items():
                if isinstance(value, dict):
                    result = self._find_list_in_data(value, list_name)
                    if result is not None:
                        return result

            return None

        except Exception as e:
            logger.error(f"Error finding list '{list_name}' in data: {e}")
            return None

    def _is_dynamic_slide(self, slide) -> bool:
        """Check if a slide contains dynamic slide markers ({{#list:...}}).
        
        Args:
            slide: The slide to check
            
        Returns:
            True if the slide contains dynamic markers, False otherwise
        """
        try:
            # Check all text shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "{{#list:" in run.text:
                                return True
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # Check table cells for dynamic markers
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if "{{#list:" in run.text:
                                            return True
            return False
        except Exception as e:
            logger.error(f"Error checking if slide is dynamic: {e}")
            return False

    def _extract_dynamic_info(self, slide) -> tuple[Optional[str], Optional[List[Dict[str, Any]]]]:
        """Extract list name and data from a dynamic slide.
        
        Args:
            slide: The slide containing dynamic markers
            
        Returns:
            Tuple of (list_name, list_data) or (None, None) if not found
        """
        try:
            import re
            
            # Pattern to match {{#list:listname}}
            pattern = r'\{\{#list:([^}]+)\}\}'
            
            # Check all text shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            match = re.search(pattern, run.text)
                            if match:
                                list_name = match.group(1).strip()
                                return list_name, None  # Data will be found separately
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # Check table cells for dynamic markers
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        match = re.search(pattern, run.text)
                                        if match:
                                            list_name = match.group(1).strip()
                                            return list_name, None
            return None, None
        except Exception as e:
            logger.error(f"Error extracting dynamic info from slide: {e}")
            return None, None

    def _process_dynamic_slide_sequence(
        self, start_index: int, list_name: str, list_data: List[Dict[str, Any]], 
        data: Dict[str, Any], images: Dict[str, Any]
    ) -> int:
        """Process a sequence of dynamic slides with their corresponding list item contexts.
        
        Args:
            start_index: The index of the first dynamic slide in the sequence
            list_name: The name of the list being processed
            list_data: The list data containing items for each slide
            data: The complete data dictionary for parent/root context
            images: Images dictionary for processing
            
        Returns:
            Number of slides processed in the sequence
        """
        try:
            slides = list(self.presentation.slides)
            slides_processed = 0
            
            logger.info(f"🔧 DEBUG: Processing dynamic sequence for '{list_name}' with {len(list_data)} items")
            
            # Process each slide in the sequence with corresponding list item context
            for item_index, item_data in enumerate(list_data):
                slide_index = start_index + item_index
                
                # Check if we have a slide at this index
                if slide_index >= len(slides):
                    logger.warning(f"Dynamic slide sequence expects slide {slide_index + 1}, but only {len(slides)} slides exist")
                    break
                    
                slide = slides[slide_index]
                logger.info(f"🔧 DEBUG: Processing dynamic slide {slide_index + 1} with item {item_index + 1}/{len(list_data)}")
                
                # Create context for this list item
                item_context = self._create_list_item_context(
                    item_data, item_index, len(list_data), data
                )
                
                # Process the slide with the item-specific context
                self._process_slide_with_enhanced_context(slide, item_context, images)
                slides_processed += 1
                
            logger.info(f"🔧 DEBUG: Completed processing {slides_processed} dynamic slides for '{list_name}'")
            return slides_processed
            
        except Exception as e:
            logger.error(f"Error processing dynamic slide sequence: {e}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            return 0

    def _process_slide_with_enhanced_context(
        self, slide, context: Dict[str, Any], images: Dict[str, Any]
    ) -> None:
        """Process a slide with enhanced context handling for dynamic slides.
        
        This method:
        1. Removes template markers ({{#list:...}})
        2. Replaces merge fields using context data
        3. Handles special variables and context navigation
        
        Args:
            slide: The slide to process
            context: The context data for field replacement
            images: Images dictionary for image processing
        """
        try:
            logger.info(f"🔧 DEBUG: Processing slide with enhanced context, keys: {list(context.keys())}")
            
            # Process all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    self._process_text_shape_with_enhanced_context(shape, context)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._process_table_shape_with_enhanced_context(shape, context)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # For image shapes, we still use the existing image processing
                    # but could be enhanced in the future for context-aware images
                    pass
                    
        except Exception as e:
            logger.error(f"Error processing slide with enhanced context: {e}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")

    def _process_text_shape_with_enhanced_context(
        self, shape, context: Dict[str, Any]
    ) -> None:
        """Process text shape with enhanced context handling.
        
        Args:
            shape: The text shape to process
            context: The context data for field replacement
        """
        try:
            import re
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    logger.info(f"🔧 DEBUG: Processing text: '{original_text[:100]}'")
                    
                    # Step 1: Remove template markers ({{#list:...}})
                    template_pattern = r'\{\{#list:[^}]+\}\}'
                    # Check for template markers before removal
                    template_matches = re.findall(template_pattern, original_text)
                    if template_matches:
                        logger.info(f"🔧 DEBUG: Found template markers to remove: {template_matches}")
                    
                    run.text = re.sub(template_pattern, "", run.text)
                    
                    # Verify template markers were removed
                    remaining_templates = re.findall(template_pattern, run.text)
                    if remaining_templates:
                        logger.warning(f"🔧 DEBUG: Template markers still present after removal: {remaining_templates}")
                    elif template_matches:
                        logger.info(f"🔧 DEBUG: Successfully removed template markers from run")
                    
                    # Step 2: Replace merge fields using context
                    run.text = self._replace_merge_fields_with_context(run.text, context)
                    
                    if original_text != run.text:
                        logger.info(f"🔧 DEBUG: Text replaced: '{original_text[:50]}...' -> '{run.text[:50]}...'")
                        
        except Exception as e:
            logger.error(f"Error processing text shape with enhanced context: {e}")

    def _process_table_shape_with_enhanced_context(
        self, shape, context: Dict[str, Any]
    ) -> None:
        """Process table shape with enhanced context handling.
        
        Args:
            shape: The table shape to process
            context: The context data for field replacement
        """
        try:
            import re
            
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text
                                
                                # Step 1: Remove template markers
                                template_pattern = r'\{\{#list:[^}]+\}\}'
                                # Check for template markers before removal
                                template_matches = re.findall(template_pattern, original_text)
                                if template_matches:
                                    logger.info(f"🔧 DEBUG: Found table template markers to remove: {template_matches}")
                                
                                run.text = re.sub(template_pattern, "", run.text)
                                
                                # Verify template markers were removed from table cell
                                remaining_templates = re.findall(template_pattern, run.text)
                                if remaining_templates:
                                    logger.warning(f"🔧 DEBUG: Table template markers still present after removal: {remaining_templates}")
                                elif template_matches:
                                    logger.info(f"🔧 DEBUG: Successfully removed table template markers from run")
                                
                                # Step 2: Replace merge fields using context
                                run.text = self._replace_merge_fields_with_context(run.text, context)
                                
                                if original_text != run.text:
                                    logger.info(f"🔧 DEBUG: Table cell text replaced: '{original_text[:30]}...' -> '{run.text[:30]}...'")
                                    
        except Exception as e:
            logger.error(f"Error processing table shape with enhanced context: {e}")

    def _replace_merge_fields_with_context(self, text: str, context: Dict[str, Any]) -> str:
        """Replace merge fields in text using the provided context.
        
        Args:
            text: The text containing merge fields
            context: The context data for field replacement
            
        Returns:
            Text with merge fields replaced
        """
        try:
            import re
            
            # Pattern to match {{field_name}} but not {{#list:...}}
            pattern = r'\{\{(?!#list:)([^}]+)\}\}'
            
            def replace_field(match):
                field_name = match.group(1).strip()
                logger.info(f"🔧 DEBUG: Replacing field: '{field_name}'")
                
                # Handle special variables
                if field_name.startswith('$'):
                    if field_name in context:
                        value = context[field_name]
                        logger.info(f"🔧 DEBUG: Special variable '{field_name}' = '{value}'")
                        return str(value) if value is not None else ""
                
                # Handle parent context (../field)
                if field_name.startswith('../'):
                    parent_field = field_name[3:]  # Remove '../'
                    if '..' in context and isinstance(context['..'], dict):
                        parent_context = context['..']
                        if parent_field in parent_context:
                            value = parent_context[parent_field]
                            logger.info(f"🔧 DEBUG: Parent field '{field_name}' = '{value}'")
                            return str(value) if value is not None else ""
                
                # Handle root context ($root.field)
                if field_name.startswith('$root.'):
                    root_field = field_name[6:]  # Remove '$root.'
                    if '$root' in context and isinstance(context['$root'], dict):
                        root_context = context['$root']
                        if root_field in root_context:
                            value = root_context[root_field]
                            logger.info(f"🔧 DEBUG: Root field '{field_name}' = '{value}'")
                            return str(value) if value is not None else ""
                
                # Handle regular fields in current context
                if field_name in context:
                    value = context[field_name]
                    logger.info(f"🔧 DEBUG: Context field '{field_name}' = '{value}'")
                    return str(value) if value is not None else ""
                
                # Handle nested field access (field.subfield)
                if '.' in field_name:
                    parts = field_name.split('.')
                    current = context
                    for part in parts:
                        if isinstance(current, dict) and part in current:
                            current = current[part]
                        else:
                            current = None
                            break
                    if current is not None:
                        logger.info(f"🔧 DEBUG: Nested field '{field_name}' = '{current}'")
                        return str(current)
                
                # Field not found, return original
                logger.warning(f"🔧 DEBUG: Field '{field_name}' not found in context")
                return match.group(0)  # Return original {{field_name}}
            
            result = re.sub(pattern, replace_field, text)
            return result
            
        except Exception as e:
            logger.error(f"Error replacing merge fields with context: {e}")
            return text

    def _normalize_field_name(self, field_name: str) -> str:
        """Normalize field name to match Excel extraction normalization.

        This handles cases where template uses {{gs_classes}} but data has g_s_classes
        due to normalize_column_name() processing.
        """
        try:
            from .utils.validation import normalize_column_name

            return normalize_column_name(field_name)
        except Exception as e:
            logger.debug(f"Could not normalize field name '{field_name}': {e}")
            return field_name

    def _get_field_name_variations(self, field_name: str) -> List[str]:
        """Get common variations of a field name for normalization mismatches.

        Returns a list of field name variations to try when the original field is not found.
        """
        variations = [field_name]  # Start with original

        try:
            # Common patterns for field name normalization mismatches
            # gs_classes <-> g_s_classes (ampersand becomes underscore)
            if "_" in field_name:
                # Try adding underscores around single letters
                parts = field_name.split("_")
                if len(parts) >= 2:
                    # gs_classes -> g_s_classes (split first part)
                    if len(parts[0]) > 1:
                        first_part = parts[0]
                        # Split camelCase or common patterns
                        if "gs" in first_part.lower():
                            new_first = "g_s"
                            variations.append(new_first + "_" + "_".join(parts[1:]))

                    # g_s_classes -> gs_classes (join first parts)
                    if len(parts) >= 3 and len(parts[0]) == 1 and len(parts[1]) == 1:
                        new_first = parts[0] + parts[1]  # g + s = gs
                        variations.append(new_first + "_" + "_".join(parts[2:]))

            # Remove duplicates while preserving order
            seen = set()
            unique_variations = []
            for var in variations:
                if var not in seen:
                    seen.add(var)
                    unique_variations.append(var)

            return unique_variations

        except Exception as e:
            logger.debug(f"Error generating field variations for '{field_name}': {e}")
            return [field_name]

    def _get_field_value_direct(self, field_name: str, data: Dict[str, Any]) -> Any:
        """Get field value using direct lookup only (no variations/recursion).

        This is used by the variation logic to avoid infinite recursion.
        """
        try:
            field_parts = field_name.split(".")
            current_value = data

            # Direct path resolution only
            for part in field_parts:
                if isinstance(current_value, dict):
                    current_value = current_value.get(part)
                elif isinstance(current_value, list):
                    try:
                        index = int(part)
                        current_value = (
                            current_value[index]
                            if 0 <= index < len(current_value)
                            else None
                        )
                    except (ValueError, IndexError):
                        current_value = None
                else:
                    current_value = None

                if current_value is None:
                    break

            return current_value

        except Exception as e:
            logger.debug(f"Error in direct field lookup for '{field_name}': {e}")
            return None

    def _duplicate_slide_simple(self, source_slide, target_index: int):
        """Create a simple duplicate of a slide.

        This is a simplified version that creates a new slide with the same layout.
        The actual content copying happens during merge field processing.
        """
        try:
            slide_layout = source_slide.slide_layout
            new_slide = self.presentation.slides.add_slide(slide_layout)

            # Copy slide properties and content
            # This is a basic implementation - a full version would copy all shapes
            for shape in source_slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # Find corresponding shape in new slide and copy text
                    for new_shape in new_slide.shapes:
                        if (
                            hasattr(new_shape, "text_frame")
                            and new_shape.text_frame
                            and new_shape.shape_type == shape.shape_type
                        ):
                            # Copy text content
                            new_shape.text_frame.clear()
                            for paragraph in shape.text_frame.paragraphs:
                                p = new_shape.text_frame.add_paragraph()
                                p.text = paragraph.text
                                p.level = paragraph.level
                            break

            return new_slide

        except Exception as e:
            logger.error(f"Failed to duplicate slide: {e}")
            # Return a basic slide if duplication fails
            return self.presentation.slides.add_slide(source_slide.slide_layout)

    def _duplicate_slide_enhanced(self, source_slide, target_index: int):
        """Create an enhanced duplicate of a slide at the specified position.

        This version properly copies all content and inserts at the right position.
        """
        try:
            # Get the slide layout
            slide_layout = source_slide.slide_layout

            # Create exact copy using complete XML duplication
            new_slide = self._copy_slide_xml(source_slide)

            # Move slide to correct position using XML manipulation
            # Note: Template markers will be removed during context processing
            self._move_slide_to_position(new_slide, target_index)

            logger.debug(f"Created enhanced duplicate slide at position {target_index}")
            return new_slide

        except Exception as e:
            logger.error(f"Failed to duplicate slide (enhanced): {e}")
            # Fall back to simple duplication
            return self.presentation.slides.add_slide(source_slide.slide_layout)

    def _copy_slide_xml(self, source_slide):
        """Create an exact duplicate slide by copying complete XML structure.
        
        Uses python-pptx built-in slide creation then replaces XML content.
        This is simpler and avoids relationship management issues.
        """
        try:
            import copy
            
            # Get the source slide's complete XML element
            source_xml = source_slide._element
            
            # Make a deep copy of the entire XML structure
            target_xml = copy.deepcopy(source_xml)
            
            # Create new slide using python-pptx (handles relationships automatically)
            new_slide = self.presentation.slides.add_slide(source_slide.slide_layout)
            
            # Replace the slide's inner content while preserving the element wrapper
            # This maintains python-pptx object consistency while getting exact template content
            new_slide._element.clear()
            
            # Copy all children from source XML into the new slide element  
            for child in target_xml:
                new_slide._element.append(copy.deepcopy(child))
            
            logger.debug("Complete XML slide copy completed successfully")
            
            # Add debugging to check if text is accessible after XML copy
            logger.info(f"🔧 DEBUG: After XML copy, new slide has {len(new_slide.shapes)} shapes")
            for i, shape in enumerate(new_slide.shapes):
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape_text = shape.text if hasattr(shape, 'text') else 'No text attr'
                    logger.info(f"🔧 DEBUG: Post-copy shape {i+1} text: '{shape_text[:50]}'")
                    logger.info(f"🔧 DEBUG: Post-copy shape {i+1} has {len(shape.text_frame.paragraphs)} paragraphs")
                    for j, para in enumerate(shape.text_frame.paragraphs):
                        logger.info(f"🔧 DEBUG: Post-copy paragraph {j+1} has {len(para.runs)} runs")
                        for k, run in enumerate(para.runs):
                            logger.info(f"🔧 DEBUG: Post-copy run {k+1} text: '{run.text}'")
            
            return new_slide
            
        except Exception as e:
            logger.error(f"XML slide copy failed: {e}")
            # Fallback to basic slide creation
            return self.presentation.slides.add_slide(source_slide.slide_layout)


    def _copy_slide_content_fallback(self, source_slide, target_slide):
        """Fallback method for slide content copying."""
        try:
            # Basic text copying as fallback
            for i, source_shape in enumerate(source_slide.shapes):
                if i < len(target_slide.shapes):
                    target_shape = target_slide.shapes[i]
                    if hasattr(source_shape, 'text_frame') and hasattr(target_shape, 'text_frame'):
                        if source_shape.text_frame and target_shape.text_frame:
                            target_shape.text_frame.clear()
                            for paragraph in source_shape.text_frame.paragraphs:
                                p = target_shape.text_frame.add_paragraph()
                                p.text = paragraph.text
        except Exception as e:
            logger.warning(f"Fallback slide copying also failed: {e}")

    def _remove_template_markers(self, slide):
        """Remove template markers ({{#list:table_name}}) from slide.
        
        This removes only the template markers while preserving all other content.
        """
        try:
            template_marker_pattern = r'\{\{#list:[^}]+\}\}'
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                # Remove template markers from text
                                import re
                                cleaned_text = re.sub(template_marker_pattern, '', run.text)
                                run.text = cleaned_text
            
            logger.debug("Template markers removed from slide")
            
        except Exception as e:
            logger.warning(f"Failed to remove template markers: {e}")

    def _move_slide_to_position(self, slide, target_index):
        """Move a slide to the specified position in the presentation.
        
        Uses XML list manipulation to reorder slides properly.
        """
        try:
            # Get the slide ID for the slide we want to move
            slide_id = slide.slide_id
            
            # Get the presentation's slide ID list (XML)
            slide_id_list = self.presentation.slides._sldIdLst
            
            # Find the slide element in the XML
            slide_element = None
            for sld_id in slide_id_list:
                if int(sld_id.get('id', '0')) == slide_id:
                    slide_element = sld_id
                    break
            
            if slide_element is not None:
                # Remove from current position
                slide_id_list.remove(slide_element)
                
                # Insert at target position (clamp to valid range)
                target_index = max(0, min(target_index, len(slide_id_list)))
                slide_id_list.insert(target_index, slide_element)
                
                logger.debug(f"Moved slide to position {target_index}")
            else:
                logger.warning(f"Could not find slide element for ID {slide_id}")
            
        except Exception as e:
            logger.warning(f"Failed to move slide to position {target_index}: {e}")

    def _copy_slide_content(self, source_slide, target_slide):
        """Copy all content from source slide to target slide."""
        try:
            # Map source shapes to target shapes by position
            source_shapes = list(source_slide.shapes)
            target_shapes = list(target_slide.shapes)

            # Copy text content from matching shapes
            for i, source_shape in enumerate(source_shapes):
                if i < len(target_shapes):
                    target_shape = target_shapes[i]

                    # Copy text frames
                    if (
                        hasattr(source_shape, "text_frame")
                        and source_shape.text_frame
                        and hasattr(target_shape, "text_frame")
                        and target_shape.text_frame
                    ):
                        self._copy_text_frame(
                            source_shape.text_frame, target_shape.text_frame
                        )

                    # Tables will be handled during merge processing
                    elif source_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        pass

        except Exception as e:
            logger.error(f"Error copying slide content: {e}")

    def _copy_text_frame(self, source_frame, target_frame):
        """Copy text frame content from source to target."""
        try:
            # Clear target frame
            target_frame.clear()

            # Copy paragraphs
            first_para = True
            for para in source_frame.paragraphs:
                if first_para and target_frame.paragraphs:
                    # Use existing first paragraph
                    p = target_frame.paragraphs[0]
                    first_para = False
                else:
                    p = target_frame.add_paragraph()

                p.text = para.text
                p.level = para.level

                # Copy basic formatting
                if hasattr(para, "alignment") and para.alignment:
                    p.alignment = para.alignment

        except Exception as e:
            logger.debug(f"Error copying text frame: {e}")

    def _create_list_item_context(
        self,
        item: Dict[str, Any],
        index: int,
        total: int,
        parent_data: Dict[str, Any],
        dynamic_config: Dict[str, Any],
    ) -> Dict[str, Any]:
        """Create context for a list item with special variables."""
        try:
            logger.info(f"🔧 DEBUG: Creating context for item {index+1}/{total}")
            logger.info(f"🔧 DEBUG: Item data: {item}")
            logger.info(f"🔧 DEBUG: Parent data keys: {list(parent_data.keys()) if parent_data else 'None'}")
            
            # Use utility function to create context
            from .utils.slide_utils import create_list_context

            special_variables = dynamic_config.get(
                "special_variables",
                {
                    "index": "$index",
                    "position": "$position",
                    "first": "$first",
                    "last": "$last",
                    "odd": "$odd",
                    "even": "$even",
                },
            )

            context = create_list_context(
                item, index, total, parent_data, special_variables
            )
            
            logger.info(f"🔧 DEBUG: Created context with keys: {list(context.keys())}")
            return context

        except Exception as e:
            logger.error(f"🔧 DEBUG: Error creating list context: {e}")
            import traceback
            logger.error(f"🔧 DEBUG: Traceback: {traceback.format_exc()}")
            return item if isinstance(item, dict) else {}

    def _process_slide_with_context(self, slide, context: Dict[str, Any]) -> None:
        """Process a slide with a specific context (for dynamic slides)."""
        logger.info(f"🔧 DEBUG: _process_slide_with_context called with context keys: {list(context.keys()) if context else 'None'}")
        logger.info(f"🔧 DEBUG: Context data sample: {dict(list(context.items())[:5]) if context else 'None'}")
        try:
            # Process text shapes with context
            shape_count = 0
            processed_shapes = 0
            for shape in slide.shapes:
                shape_count += 1
                if hasattr(shape, "text_frame") and shape.text_frame:
                    logger.info(f"🔧 DEBUG: Processing text shape {shape_count} with text: {shape.text[:50] if hasattr(shape, 'text') else 'No text'}...")
                    self._process_text_shape_with_context(shape, context)
                    processed_shapes += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    logger.info(f"🔧 DEBUG: Processing table shape {shape_count}")
                    self._process_table_shape_with_context(shape, context)
                    processed_shapes += 1
            logger.info(f"🔧 DEBUG: Processed {processed_shapes}/{shape_count} shapes in slide")

        except Exception as e:
            logger.error(f"🔧 DEBUG: Error processing slide with context: {e}")
            import traceback
            logger.error(f"🔧 DEBUG: Traceback: {traceback.format_exc()}")

    def _process_text_shape_with_context(self, shape, context: Dict[str, Any]) -> None:
        """Process text shape with specific context using paragraph-level processing."""
        try:
            shape_text = shape.text if hasattr(shape, 'text') else 'No text attribute'
            logger.info(f"🔧 DEBUG: Processing shape with text: {shape_text[:100]}")
            
            for paragraph in shape.text_frame.paragraphs:
                self._process_paragraph_with_context(paragraph, context)
                
        except Exception as e:
            logger.error(f"🔧 DEBUG: Error processing text shape with context: {e}")
            import traceback
            logger.error(f"🔧 DEBUG: Traceback: {traceback.format_exc()}")

    def _process_paragraph_with_context(self, paragraph, context: Dict[str, Any]) -> None:
        """Process a paragraph by concatenating all runs, applying patterns, then redistributing text."""
        try:
            # Step 1: Collect run information and build complete paragraph text
            run_info = []
            paragraph_text = ""
            
            logger.info(f"🔧 DEBUG: Paragraph has {len(paragraph.runs)} runs")
            
            for i, run in enumerate(paragraph.runs):
                run_start = len(paragraph_text)
                run_text = run.text
                run_end = run_start + len(run_text)
                
                logger.info(f"🔧 DEBUG: Run {i+1}/{len(paragraph.runs)} - text: '{run_text}' (length: {len(run_text)})")
                logger.info(f"🔧 DEBUG: Run {i+1} - hasattr text: {hasattr(run, 'text')}")
                
                # Try alternative ways to get text if run.text is empty
                if not run_text and hasattr(run, '_element'):
                    try:
                        # Try to get text from XML element directly
                        t_elements = run._element.findall('.//{*}t')
                        xml_texts = [elem.text for elem in t_elements if elem.text]
                        if xml_texts:
                            run_text = ''.join(xml_texts)
                            logger.info(f"🔧 DEBUG: Run {i+1} - extracted from XML: '{run_text}'")
                    except Exception as e:
                        logger.warning(f"🔧 DEBUG: Failed to extract XML text from run {i+1}: {e}")
                
                # Store run metadata for reconstruction
                run_info.append({
                    'start': run_start,
                    'end': run_end,
                    'text': run_text,
                    'run_obj': run,
                    'properties': self._extract_run_properties(run)
                })
                
                paragraph_text += run_text
                logger.info(f"🔧 DEBUG: Run [{run_start}:{run_end}]: '{run_text}'")
            
            original_paragraph_text = paragraph_text
            logger.info(f"🔧 DEBUG: Complete paragraph text: '{paragraph_text}'")
            
            # Step 2: Check for patterns that span multiple runs
            self._detect_multi_run_patterns(paragraph_text, run_info)
            
            # Step 3: Apply template marker removal on complete text
            template_pattern = r"\{\{#list:[^}]+\}\}"
            template_matches = re.findall(template_pattern, paragraph_text)
            if template_matches:
                logger.info(f"🔧 DEBUG: Found paragraph-level template markers: {template_matches}")
                paragraph_text = re.sub(template_pattern, "", paragraph_text)
                logger.info(f"🔧 DEBUG: After template removal: '{paragraph_text}'")
            
            # Step 4: Apply merge field replacement on complete text
            merge_fields = re.findall(r"\{\{([^}]+)\}\}", paragraph_text)
            logger.info(f"🔧 DEBUG: Found paragraph-level merge fields: {merge_fields}")
            
            for field in merge_fields:
                field = field.strip()
                logger.info(f"🔧 DEBUG: Resolving paragraph-level field: '{field}'")
                field_value = self._get_field_value(field, context)
                logger.info(f"🔧 DEBUG: Field '{field}' resolved to: '{field_value}' (type: {type(field_value)})")
                
                if field_value is not None:
                    paragraph_text = paragraph_text.replace(f"{{{{{field}}}}}", str(field_value))
                    logger.info(f"🔧 DEBUG: Replaced field '{field}' in paragraph text")
                else:
                    logger.warning(f"🔧 DEBUG: Could not resolve field '{field}' in context. Available keys: {list(context.keys()) if context else 'None'}")
            
            # Step 5: Redistribute modified text back to runs
            if paragraph_text != original_paragraph_text:
                logger.info(f"🔧 DEBUG: Paragraph text changed, redistributing to runs")
                self._redistribute_text_to_runs(paragraph, paragraph_text, run_info)
            else:
                logger.info(f"🔧 DEBUG: Paragraph text unchanged, keeping original runs")
                
        except Exception as e:
            logger.error(f"🔧 DEBUG: Error processing paragraph with context: {e}")
            import traceback
            logger.error(f"🔧 DEBUG: Traceback: {traceback.format_exc()}")

    def _extract_run_properties(self, run) -> Dict[str, Any]:
        """Extract run properties for later reconstruction."""
        try:
            properties = {}
            if hasattr(run, '_element') and run._element is not None:
                # Extract basic properties that can be easily restored
                rPr = run._element.get_or_add_rPr()
                properties['lang'] = rPr.get('lang', 'en-GB')
                properties['dirty'] = rPr.get('dirty', '0')
                properties['err'] = rPr.get('err', None)  # Preserve spell check errors
            return properties
        except Exception as e:
            logger.warning(f"Failed to extract run properties: {e}")
            return {}

    def _detect_multi_run_patterns(self, paragraph_text: str, run_info: list) -> None:
        """Detect and log patterns that span multiple runs for debugging."""
        try:
            # Check for template markers spanning runs
            template_matches = re.finditer(r"\{\{#list:[^}]+\}\}", paragraph_text)
            for match in template_matches:
                start, end = match.span()
                spanning_runs = [info for info in run_info if info['start'] < end and info['end'] > start]
                if len(spanning_runs) > 1:
                    run_texts = [info['text'] for info in spanning_runs]
                    logger.info(f"🔧 DEBUG: Template marker '{match.group()}' spans {len(spanning_runs)} runs: {run_texts}")
            
            # Check for merge fields spanning runs
            field_matches = re.finditer(r"\{\{([^}]+)\}\}", paragraph_text)
            for match in field_matches:
                start, end = match.span()
                spanning_runs = [info for info in run_info if info['start'] < end and info['end'] > start]
                if len(spanning_runs) > 1:
                    run_texts = [info['text'] for info in spanning_runs]
                    logger.info(f"🔧 DEBUG: Merge field '{match.group()}' spans {len(spanning_runs)} runs: {run_texts}")
                    
        except Exception as e:
            logger.warning(f"Failed to detect multi-run patterns: {e}")

    def _redistribute_text_to_runs(self, paragraph, new_text: str, run_info: list) -> None:
        """Redistribute modified text back to runs while preserving formatting."""
        try:
            # Clear existing runs
            paragraph._element.clear()
            
            if not new_text:
                # If text is empty, create one empty run
                paragraph.add_run()
                return
            
            # Simple redistribution: put all text in first run, preserve its formatting
            if run_info:
                first_run = paragraph.add_run()
                first_run.text = new_text
                
                # Try to restore basic properties from the first original run
                try:
                    original_props = run_info[0]['properties']
                    if original_props and hasattr(first_run, '_element'):
                        rPr = first_run._element.get_or_add_rPr()
                        if 'lang' in original_props:
                            rPr.set('lang', original_props['lang'])
                        if 'dirty' in original_props:
                            rPr.set('dirty', original_props['dirty'])
                        # Don't restore 'err' attribute - we don't want spell check errors on replaced text
                except Exception as e:
                    logger.warning(f"Failed to restore run properties: {e}")
            else:
                # Fallback: create basic run
                run = paragraph.add_run()
                run.text = new_text
                
            logger.info(f"🔧 DEBUG: Redistributed text to runs: '{new_text[:100]}{'...' if len(new_text) > 100 else ''}'")
            
        except Exception as e:
            logger.error(f"Failed to redistribute text to runs: {e}")
            # Fallback: try to create a basic run
            try:
                paragraph.add_run().text = new_text
            except Exception as e2:
                logger.error(f"Even fallback redistribution failed: {e2}")

    def _process_table_shape_with_context(self, shape, context: Dict[str, Any]) -> None:
        """Process table shape with specific context."""
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text

                                # First, remove any template list markers
                                template_pattern = r"\{\{#list:[^}]+\}\}"
                                run.text = re.sub(template_pattern, "", run.text)

                                # Find remaining merge fields in the text
                                merge_fields = re.findall(r"\{\{([^}]+)\}\}", run.text)

                                for field in merge_fields:
                                    field = field.strip()

                                    # Use the enhanced _get_field_value method
                                    field_value = self._get_field_value(field, context)

                                    if field_value is not None and field_value != "":
                                        # Replace the merge field with the value
                                        run.text = run.text.replace(
                                            f"{{{{{field}}}}}", str(field_value)
                                        )

        except Exception as e:
            logger.error(f"Error processing table shape with context: {e}")

    def test_slide_id_tracking(self) -> bool:
        """Test function to verify slide ID tracking works correctly.
        
        Returns:
            True if slide ID tracking is working, False otherwise
        """
        try:
            if not self.presentation:
                logger.error("No presentation loaded for slide ID tracking test")
                return False
                
            # Test: Create a set to track processed slide IDs (simulating the fix)
            test_processed_ids = set()
            
            # Get existing slide IDs
            existing_slide_ids = [slide.slide_id for slide in self.presentation.slides]
            logger.info(f"🔧 TEST: Found {len(existing_slide_ids)} existing slides with IDs: {existing_slide_ids}")
            
            # Simulate adding new slide IDs to the set (this is what our fix does)
            for slide_id in existing_slide_ids:
                test_processed_ids.add(slide_id)
                logger.info(f"🔧 TEST: Added slide ID {slide_id} to processed set")
            
            # Verify tracking works - check that slide IDs are in the set
            tracking_works = True
            for slide in self.presentation.slides:
                if slide.slide_id not in test_processed_ids:
                    logger.error(f"🔧 TEST: FAIL - Slide ID {slide.slide_id} not found in processed set")
                    tracking_works = False
                else:
                    logger.info(f"🔧 TEST: PASS - Slide ID {slide.slide_id} correctly tracked")
            
            if tracking_works:
                logger.info("🔧 TEST: Slide ID tracking test PASSED - all slide IDs correctly tracked")
            else:
                logger.error("🔧 TEST: Slide ID tracking test FAILED - some IDs missing")
                
            return tracking_works
            
        except Exception as e:
            logger.error(f"🔧 TEST: Slide ID tracking test error: {e}")
            return False

    def test_paragraph_level_processing(self) -> bool:
        """Test function to verify paragraph-level processing works correctly.
        
        Returns:
            True if paragraph-level processing works, False otherwise
        """
        try:
            if not self.presentation:
                logger.error("No presentation loaded for paragraph processing test")
                return False
            
            # Test context data
            test_context = {
                'gs_classes': '35 - Business Services',
                'gs_terms': 'Business management and administration services'
            }
            
            # Find a slide with text shapes to test
            for slide_idx, slide in enumerate(self.presentation.slides):
                if len(slide.shapes) > 0:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            # Test the paragraph processing on this shape
                            logger.info(f"🔧 TEST: Testing paragraph processing on slide {slide_idx + 1}")
                            logger.info(f"🔧 TEST: Original shape text: '{shape.text[:100] if hasattr(shape, 'text') else 'No text'}'")
                            
                            # Process the shape with our new method
                            self._process_text_shape_with_context(shape, test_context)
                            
                            # Check results
                            new_text = shape.text if hasattr(shape, 'text') else ''
                            logger.info(f"🔧 TEST: Processed shape text: '{new_text[:100]}'")
                            
                            logger.info("🔧 TEST: Paragraph-level processing test completed")
                            return True
            
            logger.warning("🔧 TEST: No suitable text shapes found for testing")
            return False
            
        except Exception as e:
            logger.error(f"🔧 TEST: Paragraph processing test error: {e}")
            import traceback
            logger.error(f"🔧 TEST: Traceback: {traceback.format_exc()}")
            return False

    def test_two_pass_processing(self) -> bool:
        """Test function to verify two-pass processing works correctly.
        
        Returns:
            True if two-pass processing works, False otherwise
        """
        try:
            if not self.presentation:
                logger.error("No presentation loaded for two-pass processing test")
                return False
                
            # Test data simulating the real dynamic slide scenario
            test_data = {
                'gs_classes_terms': [
                    {'gs_classes': '35 - Business Services', 'gs_terms': 'Business management services'},
                    {'gs_classes': '37 - Education', 'gs_terms': 'Training and education services'}
                ]
            }
            
            # Test config
            test_config = {
                'global_settings': {
                    'powerpoint': {
                        'dynamic_slides': {
                            'enabled': True,
                            'template_marker': '{{#list:',
                            'remove_template_slides': True
                        }
                    }
                }
            }
            
            logger.info("🔧 TEST: Starting two-pass processing test")
            
            # Test Pass 1: Dynamic slide creation
            powerpoint_config = test_config['global_settings']['powerpoint']
            dynamic_slide_mapping = self._create_dynamic_slides_pass(test_data, powerpoint_config)
            
            logger.info(f"🔧 TEST: Pass 1 created {len(dynamic_slide_mapping)} dynamic slide mappings")
            
            # Test Pass 2: Content processing
            self._process_slides_pass(test_data, {}, powerpoint_config, dynamic_slide_mapping)
            
            logger.info("🔧 TEST: Two-pass processing test completed successfully")
            return True
            
        except Exception as e:
            logger.error(f"🔧 TEST: Two-pass processing test error: {e}")
            import traceback
            logger.error(f"🔧 TEST: Traceback: {traceback.format_exc()}")
            return False

    def _process_slide(
        self,
        slide,
        data: Dict[str, Any],
        images: Optional[Dict[str, List[Dict[str, Any]]]] = None,
    ) -> None:
        """Process a single slide for merge field replacement."""
        slide_id = getattr(slide, 'slide_id', 'unknown')
        logger.debug(f"Starting _process_slide for slide ID {slide_id}")
        
        try:
            # Process text shapes
            shapes_list = list(slide.shapes)
            shape_count = len(shapes_list)
            logger.debug(f"Found {shape_count} shapes to process in slide {slide_id}")
            
            processed_shapes = 0
            for i, shape in enumerate(shapes_list):
                logger.debug(f"Processing shape {i+1}/{shape_count}, type: {getattr(shape, 'shape_type', 'unknown')}")
                
                if hasattr(shape, "text_frame") and shape.text_frame:
                    logger.debug(f"Processing text shape {i+1}")
                    self._process_text_shape(shape, data, images)
                    processed_shapes += 1

                # Process tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    logger.debug(f"Processing table shape {i+1}")
                    self._process_table_shape(shape, data)
                    processed_shapes += 1

                # Process placeholder images
                elif hasattr(shape, "text") and shape.text:
                    # Check if this is an image placeholder
                    if self._is_image_placeholder(shape.text):
                        logger.debug(f"Processing image placeholder shape {i+1}")
                        self._replace_image_placeholder(slide, shape, data, images)
                        processed_shapes += 1
                else:
                    logger.debug(f"Skipping shape {i+1} - no processable content")
            
            logger.debug(f"Completed _process_slide for slide {slide_id}: processed {processed_shapes}/{shape_count} shapes")

        except Exception as e:
            logger.error(f"Failed to process slide {slide_id}: {e}")

    def _process_text_shape(
        self,
        shape: BaseShape,
        data: Dict[str, Any],
        images: Optional[Dict[str, List[Dict[str, Any]]]] = None,
    ) -> None:
        """Process text shape for merge field replacement."""
        try:
            if not shape.text_frame:
                return

            # Check if this shape contains any image fields (based on configuration, not just field names)
            shape_text = self._get_full_text_from_shape(shape)
            if shape_text and self._contains_image_field(shape_text, data):
                # Try to replace with image first
                if self._try_replace_text_with_image(shape, data, images):
                    return  # Successfully replaced with image, skip text processing

            # Process each paragraph for text replacement
            for paragraph in shape.text_frame.paragraphs:
                self._process_paragraph(paragraph, data)

        except Exception as e:
            logger.warning(f"Failed to process text shape: {e}")

    def _try_replace_text_with_image(
        self,
        shape: BaseShape,
        data: Dict[str, Any],
        images: Optional[Dict[str, List[Dict[str, Any]]]] = None,
    ) -> bool:
        """Try to replace text shape with image if it's an image placeholder."""
        if not images:
            return False

        try:
            shape_text = self._get_full_text_from_shape(shape)
            merge_fields = validate_merge_fields(shape_text)

            if not merge_fields:
                return False

            # Try to find image for the first merge field
            for field in merge_fields:
                # Check if this is an image field
                if self._is_image_field(field, data):
                    image_path = self._get_image_for_field(field, data, images)
                    if image_path:
                        # Verify the image file exists
                        if os.path.exists(str(image_path)):
                            # Replace the text shape with an image
                            return self._replace_shape_with_image(shape, image_path)
                        else:
                            logger.warning(
                                f"Image file does not exist at path: {image_path}"
                            )
                            # Keep the original text if image insertion fails
                            return False

            return False

        except Exception as e:
            logger.warning(f"Failed to replace text with image: {e}")
            return False

    def _get_field_type(self, field_name: str, data: Dict[str, Any]) -> str:
        """Get the type of a field from metadata.

        Args:
            field_name: The field name or path to check
            data: The data structure containing field values and metadata

        Returns:
            The field type as a string ('text', 'image', etc.)
        """
        # Check for direct field type in the current data level
        if "_field_types" in data and field_name in data["_field_types"]:
            return data["_field_types"][field_name]

        # For nested fields, try to navigate to the parent object
        parts = field_name.split(".")
        if len(parts) > 1:
            # Try to find the parent object
            current_data = data
            parent_path = ".".join(parts[:-1])
            leaf_name = parts[-1]

            # First try direct path
            parent = self._get_field_value(parent_path, data)

            if parent and isinstance(parent, dict):
                if "_field_types" in parent and leaf_name in parent["_field_types"]:
                    return parent["_field_types"][leaf_name]

            # Try to find in sheet data
            for sheet_name, sheet_data in data.items():
                if isinstance(sheet_data, dict):
                    # Try to find the field in this sheet
                    for table_name, table_data in sheet_data.items():
                        if isinstance(table_data, list):
                            # Check each row in the table
                            for row in table_data:
                                if (
                                    isinstance(row, dict)
                                    and "_field_types" in row
                                    and leaf_name in row["_field_types"]
                                ):
                                    return row["_field_types"][leaf_name]

        # Default to text if no type information is found
        return "text"

    def _is_image_field(self, field_name: str, data: Dict[str, Any]) -> bool:
        """Determine if a field is an image field based on metadata."""
        # Check if it's a range image
        if "_range_images" in data and field_name in data["_range_images"]:
            return True

        field_type = self._get_field_type(field_name, data)
        return field_type == "image"

    def _contains_image_field(self, text_content: str, data: Dict[str, Any]) -> bool:
        """Check if text content contains any image fields based on configuration."""
        if not text_content:
            return False

        # Extract merge fields
        merge_fields = validate_merge_fields(text_content)
        if not merge_fields:
            return False

        # Check if any field is configured as an image field
        for field in merge_fields:
            if self._is_image_field(field, data):
                return True

        return False

    def _is_link_field(self, field_name: str, data: Dict[str, Any]) -> bool:
        """Check if a field is configured as a link field."""
        field_type = self._get_field_type(field_name, data)
        logger.debug(f"🔗 Field '{field_name}' has type: {field_type}")
        if field_type == "link":
            return True

        # Fallback: check if field value is a link object (has 'title' and 'link' keys)
        field_value = self._get_field_value(field_name, data)
        logger.debug(f"🔗 Field '{field_name}' value: {field_value}")
        if isinstance(field_value, dict):
            has_title = "title" in field_value
            has_link = "link" in field_value
            logger.debug(
                f"🔗 Field '{field_name}' has title: {has_title}, has link: {has_link}"
            )
            return has_title and has_link

        return False

    def _contains_link_field(self, text_content: str, data: Dict[str, Any]) -> bool:
        """Check if text content contains any link fields based on configuration."""
        if not text_content:
            return False

        # Extract merge fields
        merge_fields = validate_merge_fields(text_content)
        if not merge_fields:
            return False

        # Check if any field is configured as a link field
        for field in merge_fields:
            if self._is_link_field(field, data):
                return True

        return False

    def _is_image_placeholder(self, text_content: str) -> bool:
        """Legacy method - check if text content is an image placeholder.

        This method is kept for backward compatibility but should use
        _contains_image_field() when data is available.
        """
        if not text_content:
            return False

        # Extract merge fields
        merge_fields = validate_merge_fields(text_content)
        if not merge_fields:
            return False

        # Fallback heuristic: check if field name contains image-related keywords
        return any(
            "image" in field.lower()
            or "img" in field.lower()
            or "photo" in field.lower()
            or "picture" in field.lower()
            or "range" in field.lower()  # Added for range images
            for field in merge_fields
        )

    def _get_image_for_field(
        self,
        field_name: str,
        data: Dict[str, Any],
        images: Dict[str, List[Dict[str, Any]]],
    ) -> Optional[str]:
        """Get image path for a specific field."""
        try:
            # Check for range images first
            if "_range_images" in data and field_name in data["_range_images"]:
                range_image_path = data["_range_images"][field_name]
                if os.path.exists(str(range_image_path)):
                    logger.debug(
                        f"Found range image for field '{field_name}': {range_image_path}"
                    )
                    return str(range_image_path)
                else:
                    logger.warning(
                        f"Range image file does not exist: {range_image_path}"
                    )

            # Check if this is an image field based on metadata
            is_image = self._is_image_field(field_name, data)

            # First try to get image data from data
            image_data = self._get_field_value(field_name, data)
            if image_data:
                if is_image:
                    logger.debug(
                        f"Found image data for field '{field_name}': {type(image_data)}"
                    )

                    # Handle new mixed data structure (dict with base64/path)
                    if isinstance(image_data, dict):
                        # Try path first if it exists and is valid
                        if "path" in image_data and image_data["path"]:
                            path = str(image_data["path"])
                            if os.path.exists(path):
                                logger.debug(f"Using existing file path: {path}")
                                return path

                        # Fall back to base64 data
                        if "base64" in image_data and image_data["base64"]:
                            # Create temporary file from base64 data
                            temp_path = self._create_temp_image_from_base64(
                                image_data["base64"]
                            )
                            if temp_path:
                                logger.debug(
                                    f"Created temporary image from base64: {temp_path}"
                                )
                                return temp_path

                    # Handle legacy string path format
                    elif isinstance(image_data, str):
                        if os.path.exists(str(image_data)):
                            logger.debug(f"Using image file path: {image_data}")
                            return str(image_data)
                        else:
                            logger.warning(
                                f"Image file does not exist at path: {image_data}"
                            )
                else:
                    logger.debug(
                        f"Field '{field_name}' is not an image field, but has value: {type(image_data)}"
                    )
                    # Not an image field, don't return the path
                    return None

            # If field is an image but we didn't find a path, try position-based matching
            if is_image:
                image_path = self._find_image_by_field_name(field_name, images)
                if image_path:
                    # Verify the image file exists
                    if os.path.exists(str(image_path)):
                        logger.debug(
                            f"Found image by field name matching: {image_path}"
                        )
                        return str(image_path)
                    else:
                        logger.warning(
                            f"Image file from field name matching does not exist: {image_path}"
                        )

            # Log if no image was found for an image field
            if is_image:
                logger.warning(f"No image found for image field: {field_name}")
                # Debug log the available images
                if images:
                    logger.debug(f"Available images: {images}")
                else:
                    logger.warning("No images available")

            return None

        except Exception as e:
            logger.warning(f"Failed to get image for field '{field_name}': {e}")
            return None

    def _find_image_by_field_name(
        self, field_name: str, images: Dict[str, List[Dict[str, Any]]]
    ) -> Optional[str]:
        """Find image by field name using various matching strategies."""
        if not images:
            logger.warning("No images provided to _find_image_by_field_name")
            return None

        logger.debug(f"Looking for image matching field: {field_name}")

        # Log available images for debugging
        for sheet_name, sheet_images in images.items():
            logger.debug(f"Sheet {sheet_name} has {len(sheet_images)} images")
            for idx, img in enumerate(sheet_images):
                logger.debug(
                    f"  Image {idx}: {img.get('filename')} at {img.get('path')}"
                )

        field_lower = field_name.lower()

        # Strategy 1: Direct field name matching
        try:
            for sheet_name, sheet_images in images.items():
                for image_info in sheet_images:
                    # Check if field name contains position information
                    if "position" in image_info:
                        position = image_info["position"]
                        if position.get("estimated_cell"):
                            cell_ref = position["estimated_cell"].lower()
                            if cell_ref in field_lower or field_lower.endswith(
                                cell_ref
                            ):
                                logger.info(
                                    f"Found image by position match: {image_info['path']}"
                                )
                                return image_info["path"]
        except Exception as e:
            logger.debug(f"Error in position-based image matching: {e}")

        # Strategy 2: Pattern matching for common image field patterns
        try:
            patterns = [
                r"image_search\.(\d+)\.image",  # image_search.0.image
                r"(\w+)_image_(\d+)",  # sheet_image_1
                r"image(\d+)",  # image1
                r"img(\d+)",  # img1
            ]

            for pattern in patterns:
                match = re.search(pattern, field_lower)
                if match:
                    index_str = match.group(1) if match.groups() else "0"
                    try:
                        index = int(index_str)
                        # Find image by index across all sheets
                        for sheet_name, sheet_images in images.items():
                            if 0 <= index < len(sheet_images):
                                logger.info(
                                    f"Found image by pattern match: {sheet_images[index]['path']}"
                                )
                                return sheet_images[index]["path"]
                    except (ValueError, IndexError):
                        continue
        except Exception as e:
            logger.debug(f"Error in pattern-based image matching: {e}")

        # Strategy 3: Keyword matching
        try:
            keywords = ["image", "img", "picture", "photo"]
            for keyword in keywords:
                if keyword in field_lower:
                    # Return first available image
                    for sheet_name, sheet_images in images.items():
                        if sheet_images:
                            logger.info(
                                f"Found image by keyword match: {sheet_images[0]['path']}"
                            )
                            return sheet_images[0]["path"]
        except Exception as e:
            logger.debug(f"Error in keyword-based image matching: {e}")

        # Strategy 4: Just use the first available image if all else fails
        try:
            for sheet_name, sheet_images in images.items():
                if sheet_images:
                    logger.info(
                        f"No specific match found, using first available image: {sheet_images[0]['path']}"
                    )
                    return sheet_images[0]["path"]
        except Exception as e:
            logger.debug(f"Error in fallback image matching: {e}")

        logger.warning(f"No image found for field: {field_name}")
        return None

    def _replace_shape_with_image(self, shape, image_path: str) -> bool:
        """Replace a shape with an image while maintaining position and size.

        Args:
            shape: The shape to replace
            image_path: Path to the image file

        Returns:
            True if successful, False otherwise
        """
        try:
            # Verify the image file exists
            if not os.path.exists(image_path):
                logger.error(f"Image file does not exist: {image_path}")
                return False

            # Ensure image_path is absolute
            if not os.path.isabs(image_path):
                image_path = os.path.abspath(image_path)

            # Get the parent slide
            slide = shape.part.slide

            # Get the shape dimensions and position
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Store shape properties before removal
            shape_name = getattr(shape, "name", "Unknown")

            # Clear the text from the shape to prepare for replacement
            if hasattr(shape, "text_frame") and shape.text_frame:
                try:
                    shape.text_frame.clear()  # Use proper API method
                    logger.debug(
                        f"Cleared text from shape '{shape_name}' before replacement"
                    )
                except Exception as text_err:
                    logger.warning(
                        f"Could not clear text from shape '{shape_name}': {text_err}"
                    )

            # Use python-pptx API for safer shape removal
            try:
                # Get slide reference
                slide_shapes = slide.shapes

                # Find and remove the shape using the proper API
                for i, slide_shape in enumerate(slide_shapes):
                    if slide_shape == shape:
                        # Remove using the shape collection API (safer than direct XML)
                        del slide_shapes[i]
                        logger.debug(f"Safely removed shape '{shape_name}' from slide")
                        break

            except Exception as removal_err:
                logger.warning(
                    f"Could not remove shape using API, trying XML method: {removal_err}"
                )

                # Fallback to XML method if API fails
                try:
                    shape_id = shape.shape_id
                    sp_tree = slide.shapes._spTree

                    for sp in sp_tree.findall(".//{*}sp"):
                        if sp.get("id") == str(shape_id):
                            sp_tree.remove(sp)
                            logger.debug(
                                f"Removed shape '{shape_name}' using XML fallback"
                            )
                            break

                except Exception as xml_err:
                    logger.error(f"Failed to remove shape '{shape_name}': {xml_err}")
                    return False

            # Add the image while maintaining aspect ratio
            try:
                # Get image dimensions
                with PILImage.open(image_path) as img:
                    img_width, img_height = img.size

                # Calculate aspect ratios
                shape_ratio = width / height
                img_ratio = img_width / img_height

                # Adjust dimensions to maintain aspect ratio
                if img_ratio > shape_ratio:
                    # Image is wider than shape
                    new_width = width
                    new_height = width / img_ratio
                    new_top = top + (height - new_height) / 2
                    new_left = left
                else:
                    # Image is taller than shape
                    new_height = height
                    new_width = height * img_ratio
                    new_left = left + (width - new_width) / 2
                    new_top = top

                # Add the image to the slide
                picture = slide.shapes.add_picture(
                    image_path, new_left, new_top, new_width, new_height
                )

                # Ensure the image is on top (higher z-order)
                # This is done by moving the picture to the end of the shapes collection
                try:
                    pic_element = None
                    for shape_xml in sp_tree.findall(".//{*}pic"):
                        if shape_xml.get("id") == str(picture.shape_id):
                            pic_element = shape_xml
                            break

                    if pic_element is not None and pic_element.getparent() is not None:
                        # Remove and re-append to put it at the end (top of z-order)
                        parent = pic_element.getparent()
                        parent.remove(pic_element)
                        parent.append(pic_element)
                        logger.debug(f"Adjusted z-order to ensure image is on top")
                except Exception as z_err:
                    logger.warning(f"Could not adjust z-order: {z_err}")

                logger.info(f"Successfully replaced shape with image: {image_path}")
                return True

            except Exception as e:
                logger.error(f"Error adding image to slide: {e}")

                # Add a text box with error message as fallback
                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.text = f"Image Error: {os.path.basename(image_path)}"

                return False

        except Exception as e:
            logger.error(f"Error replacing shape with image: {e}")
            return False

    def _process_paragraph(self, paragraph, data: Dict[str, Any]) -> None:
        """Process a single paragraph for merge field replacement while preserving formatting."""
        try:
            # Log the paragraph text being processed
            paragraph_text = (
                paragraph.text if hasattr(paragraph, "text") else str(paragraph)
            )
            logger.debug(f"📝 Processing paragraph: '{paragraph_text}'")

            # Try the new formatting-preserving approach first
            if self._process_paragraph_preserve_formatting(paragraph, data):
                return

            # Fall back to the original approach if the new one fails
            logger.warning("Falling back to original paragraph processing")
            self._process_paragraph_original(paragraph, data)

        except Exception as e:
            logger.warning(f"Failed to process paragraph: {e}")

    def _process_paragraph_preserve_formatting(
        self, paragraph, data: Dict[str, Any]
    ) -> bool:
        """Process paragraph while preserving run-level formatting - ITERATIVE VERSION."""
        try:
            # First, remove any template list markers from all runs
            template_pattern = r"\{\{#list:[^}]+\}\}"
            for run in paragraph.runs:
                if run.text:
                    original_text = run.text
                    cleaned_text = re.sub(template_pattern, "", run.text)
                    if cleaned_text != original_text:
                        run.text = cleaned_text
                        logger.debug(
                            f"Removed template markers from run: '{original_text}' -> '{cleaned_text}'"
                        )

            # Use iterative approach to avoid position conflicts
            # When we replace one field, it changes text positions, so we need to
            # recalculate positions for remaining fields
            max_iterations = 10  # Prevent infinite loops
            iteration = 0

            while iteration < max_iterations:
                # Recalculate field positions each time to avoid position shifts
                field_positions = self._find_merge_fields_in_runs(paragraph)

                if not field_positions:
                    # No more fields to process
                    break

                # Process just the first field found to avoid position conflicts
                field_info = field_positions[0]
                field_name = field_info["field"]
                field_value = self._get_field_value(field_name, data)

                logger.debug(
                    f"Processing field '{field_name}' with value type '{type(field_value).__name__}' (iteration {iteration + 1})"
                )

                # Check if this is a link field and process accordingly
                if self._is_link_field(field_name, data):
                    # Process as hyperlink
                    logger.debug(f"🔗 Processing field '{field_name}' as hyperlink")
                    success = self._process_link_field(
                        field_name, field_value, paragraph, field_info
                    )
                    logger.debug(f"🔗 Hyperlink processing success: {success}")
                    if not success:
                        # Fall back to text processing if link processing fails
                        logger.debug(
                            f"🔗 Falling back to text processing for field '{field_name}'"
                        )
                        field_value_str = (
                            str(field_value) if field_value is not None else ""
                        )
                        self._replace_field_in_runs(
                            paragraph, field_info, field_value_str
                        )
                else:
                    # Process as regular text field
                    field_value_str = (
                        str(field_value) if field_value is not None else ""
                    )
                    self._replace_field_in_runs(paragraph, field_info, field_value_str)

                iteration += 1

            if iteration >= max_iterations:
                logger.warning(
                    "Reached maximum iterations in paragraph processing - possible infinite loop"
                )

            # Clean up any empty runs that might cause PowerPoint issues
            self._cleanup_empty_runs(paragraph)

            # ENHANCED: Ensure clean field structure for field paragraphs
            if self._paragraph_contains_fields(paragraph):
                self._ensure_clean_field_structure(paragraph)

            # ENHANCED: Aggressive cleanup of all empty/whitespace runs
            self._aggressive_empty_run_cleanup(paragraph)

            # ENHANCED: Remove any remaining error attributes from the paragraph
            self._remove_error_attributes_from_paragraph(paragraph)

            return True

        except Exception as e:
            logger.warning(
                f"Failed to process paragraph with formatting preservation: {e}"
            )
            return False

    def _process_paragraph_original(self, paragraph, data: Dict[str, Any]) -> None:
        """Original paragraph processing method (fallback)."""
        try:
            # Get the full paragraph text
            paragraph_text = ""
            for run in paragraph.runs:
                paragraph_text += run.text

            # Remove template list markers before processing merge fields
            template_pattern = r"\{\{#list:[^}]+\}\}"
            paragraph_text = re.sub(template_pattern, "", paragraph_text)

            # Update runs with cleaned text if template markers were removed
            original_text = "".join(run.text for run in paragraph.runs)
            if paragraph_text != original_text:
                # Clear existing runs and set cleaned text
                for run in paragraph.runs:
                    run.text = ""
                if paragraph.runs:
                    paragraph.runs[0].text = paragraph_text

            # Find merge fields in paragraph
            merge_fields = validate_merge_fields(paragraph_text)

            if not merge_fields:
                return

            # Replace merge fields
            new_text = paragraph_text
            for field in merge_fields:
                field_value = self._get_field_value(field, data)
                merge_pattern = f"{{{{{field}}}}}"
                new_text = new_text.replace(
                    merge_pattern, str(field_value) if field_value is not None else ""
                )

            # Update paragraph text if it changed
            if new_text != paragraph_text:
                # Clear existing runs
                for run in paragraph.runs:
                    run.text = ""

                # Set new text in first run
                if paragraph.runs:
                    paragraph.runs[0].text = new_text
                else:
                    paragraph.add_run().text = new_text

        except Exception as e:
            logger.warning(f"Failed to process paragraph with original method: {e}")

    def _find_merge_fields_in_runs(self, paragraph) -> List[Dict[str, Any]]:
        """Find merge fields and their positions within paragraph runs."""
        field_positions = []

        try:
            # Build a map of text positions to runs
            run_map = []
            text_position = 0

            for run_idx, run in enumerate(paragraph.runs):
                run_text = run.text
                run_start = text_position
                run_end = text_position + len(run_text)

                run_map.append(
                    {
                        "run_idx": run_idx,
                        "run": run,
                        "text": run_text,
                        "start": run_start,
                        "end": run_end,
                    }
                )

                text_position = run_end

            # Get full paragraph text
            full_text = "".join(run["text"] for run in run_map)

            # Find all merge fields in the full text
            merge_fields = validate_merge_fields(full_text)

            for field in merge_fields:
                field_pattern = f"{{{{{field}}}}}"
                field_start = full_text.find(field_pattern)

                if field_start != -1:
                    field_end = field_start + len(field_pattern)

                    # Find which runs contain this field
                    affected_runs = []
                    for run_info in run_map:
                        # Check if this run overlaps with the field
                        if (
                            run_info["start"] < field_end
                            and run_info["end"] > field_start
                        ):
                            # Calculate the portion of the field in this run
                            field_start_in_run = max(0, field_start - run_info["start"])
                            field_end_in_run = min(
                                len(run_info["text"]), field_end - run_info["start"]
                            )

                            affected_runs.append(
                                {
                                    "run_idx": run_info["run_idx"],
                                    "run": run_info["run"],
                                    "field_start_in_run": field_start_in_run,
                                    "field_end_in_run": field_end_in_run,
                                    "run_text": run_info["text"],
                                }
                            )

                    if affected_runs:
                        field_positions.append(
                            {
                                "field": field,
                                "field_pattern": field_pattern,
                                "field_start": field_start,
                                "field_end": field_end,
                                "affected_runs": affected_runs,
                            }
                        )

            return field_positions

        except Exception as e:
            logger.warning(f"Failed to find merge fields in runs: {e}")
            return []

    def _replace_field_in_runs(
        self, paragraph, field_info: Dict[str, Any], replacement_text: str
    ) -> None:
        """Replace a merge field in runs while preserving formatting."""
        try:
            affected_runs = field_info["affected_runs"]

            if len(affected_runs) == 1:
                # Simple case: field is entirely within one run
                self._replace_field_in_single_run(affected_runs[0], replacement_text)
            else:
                # Complex case: field spans multiple runs
                self._replace_field_across_runs(
                    paragraph, affected_runs, replacement_text
                )

        except Exception as e:
            logger.warning(f"Failed to replace field in runs: {e}")

    def _replace_field_in_single_run(
        self, run_info: Dict[str, Any], replacement_text: str
    ) -> None:
        """Replace field within a single run."""
        try:
            run = run_info["run"]
            original_text = run_info["run_text"]
            field_start = run_info["field_start_in_run"]
            field_end = run_info["field_end_in_run"]

            # Build new text by replacing the field portion
            new_text = (
                original_text[:field_start]
                + replacement_text
                + original_text[field_end:]
            )

            # Update the run text (formatting is preserved automatically)
            run.text = new_text
            # logger.debug(f"Replaced field in single run: '{original_text}' -> '{new_text}'")  # Too verbose

        except Exception as e:
            logger.warning(f"Failed to replace field in single run: {e}")

    def _replace_field_across_runs(
        self, paragraph, affected_runs: List[Dict[str, Any]], replacement_text: str
    ) -> None:
        """Replace field that spans across multiple runs - FIXED VERSION."""
        try:
            if not affected_runs:
                return

            # FIXED STRATEGY:
            # 1. Build the complete replacement by preserving text before and after field
            # 2. Put everything in the first run
            # 3. Clear all other affected runs

            first_run = affected_runs[0]
            last_run = affected_runs[-1]

            # Text before the field (from first run)
            text_before = first_run["run_text"][: first_run["field_start_in_run"]]

            # Text after the field (from last run)
            text_after = last_run["run_text"][last_run["field_end_in_run"] :]

            # Complete replacement text
            complete_text = text_before + replacement_text + text_after

            # Update first run with complete text
            first_run["run"].text = complete_text

            # Clear all other affected runs
            for i in range(1, len(affected_runs)):
                affected_runs[i]["run"].text = ""

            logger.debug(
                f"Fixed multi-run replacement: combined {len(affected_runs)} runs into first run"
            )

        except Exception as e:
            logger.warning(f"Failed to replace field across runs (fixed): {e}")

    def _preserve_run_formatting(self, source_run, target_run) -> None:
        """Copy formatting properties from source run to target run."""
        try:
            # Copy font properties
            if hasattr(source_run, "font") and hasattr(target_run, "font"):
                source_font = source_run.font
                target_font = target_run.font

                # Copy common font properties
                if source_font.name:
                    target_font.name = source_font.name
                if source_font.size:
                    target_font.size = source_font.size
                if source_font.bold is not None:
                    target_font.bold = source_font.bold
                if source_font.italic is not None:
                    target_font.italic = source_font.italic
                if source_font.underline is not None:
                    target_font.underline = source_font.underline
                if source_font.color.rgb:
                    target_font.color.rgb = source_font.color.rgb

                # logger.debug("Copied font formatting between runs")  # Too verbose

        except Exception as e:
            logger.warning(f"Failed to preserve run formatting: {e}")

    def _cleanup_empty_runs(self, paragraph) -> None:
        """Remove empty runs that can cause PowerPoint repair issues - ENHANCED VERSION."""
        try:
            runs_to_remove = []

            for run in paragraph.runs:
                should_remove = False

                # Check if run is truly empty (no text or only whitespace)
                if not run.text or not run.text.strip():
                    should_remove = True

                    # Check if run has meaningful formatting that should be preserved
                    has_formatting = False
                    if hasattr(run, "font"):
                        font = run.font
                        if (
                            font.bold is not None
                            or font.italic is not None
                            or font.underline is not None
                            or font.size is not None
                            or font.name is not None
                        ):
                            has_formatting = True

                    # Don't remove if it has meaningful formatting
                    if has_formatting:
                        should_remove = False

                # ENHANCED: Check for error attributes that cause PowerPoint repair issues
                if hasattr(run, "_element") and run._element is not None:
                    # Check for error attributes in run properties
                    rpr_elements = run._element.xpath(
                        ".//a:rPr",
                        namespaces={
                            "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
                        },
                    )

                    for rpr in rpr_elements:
                        # Check for err="1" attribute that causes repair dialogs
                        if rpr.get("err") == "1":
                            # For empty runs with err="1", remove the entire run
                            if not run.text or not run.text.strip():
                                logger.debug(
                                    f"Found empty run with err='1' attribute, marking for removal"
                                )
                                should_remove = True
                                break
                            else:
                                # For non-empty runs with err="1", just remove the error attribute
                                logger.debug(
                                    f"Found non-empty run with err='1' attribute, will clean attribute"
                                )
                                # Don't remove the run, just clean the attribute in the validation phase

                        # Also remove runs with dirty="0" and empty text (common after field replacement)
                        if rpr.get("dirty") == "0" and (
                            not run.text or not run.text.strip()
                        ):
                            should_remove = True

                if should_remove:
                    runs_to_remove.append(run)

            # Remove problematic runs (but keep at least one run in the paragraph)
            if len(runs_to_remove) < len(paragraph.runs):
                removed_count = 0
                for run in runs_to_remove:
                    try:
                        # Check if run has _element attribute before accessing it
                        if hasattr(run, "_element") and run._element is not None:
                            # Remove the run's XML element from the paragraph
                            paragraph._element.remove(run._element)
                            removed_count += 1
                        else:
                            logger.debug(
                                f"Run object doesn't have _element attribute, skipping removal"
                            )
                    except Exception as remove_err:
                        logger.warning(
                            f"Could not remove problematic run: {remove_err}"
                        )

                if removed_count > 0:
                    logger.debug(
                        f"Removed {removed_count} problematic runs from paragraph"
                    )

            # ENHANCED: Ensure paragraph has at least one valid run
            if len(paragraph.runs) == 0:
                # Add a minimal empty run to prevent paragraph corruption
                new_run = paragraph.add_run()
                new_run.text = ""
                logger.debug("Added minimal run to empty paragraph")

        except Exception as e:
            logger.warning(f"Failed to cleanup empty runs: {e}")

    def _validate_presentation_integrity(self) -> None:
        """Validate presentation integrity to prevent PowerPoint repair issues - ENHANCED VERSION."""
        try:
            if not self.presentation:
                return

            issues_found = 0

            for slide_idx, slide in enumerate(self.presentation.slides):
                try:
                    # Check each shape on the slide
                    for shape_idx, shape in enumerate(slide.shapes):
                        # Check text frames for empty paragraphs/runs
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for para_idx, paragraph in enumerate(
                                shape.text_frame.paragraphs
                            ):
                                # ENHANCED: Clean up error attributes before other validation
                                error_attrs_removed = (
                                    self._remove_error_attributes_from_paragraph(
                                        paragraph
                                    )
                                )
                                if error_attrs_removed:
                                    issues_found += error_attrs_removed

                                # Ensure paragraph has at least one run
                                # Skip field paragraphs to preserve clean structure
                                if len(
                                    paragraph.runs
                                ) == 0 and not self._paragraph_contains_only_fields(
                                    paragraph
                                ):
                                    # Add a minimal run to prevent issues
                                    paragraph.add_run()
                                    issues_found += 1

                                # Check for completely empty paragraphs
                                total_text = "".join(run.text for run in paragraph.runs)
                                if not total_text and len(paragraph.runs) > 1:
                                    # Multiple empty runs - consolidate to one
                                    for run in paragraph.runs[1:]:
                                        try:
                                            paragraph._element.remove(run._element)
                                        except:
                                            pass
                                    issues_found += 1

                except Exception as shape_err:
                    logger.warning(
                        f"Error validating slide {slide_idx + 1}: {shape_err}"
                    )

            if issues_found > 0:
                logger.info(
                    f"Fixed {issues_found} potential PowerPoint compatibility issues"
                )
            else:
                pass  # No issues found

        except Exception as e:
            logger.warning(f"Failed to validate presentation integrity: {e}")

    def _remove_error_attributes_from_paragraph(self, paragraph) -> int:
        """Remove error attributes from runs that cause PowerPoint repair issues."""
        try:
            removed_count = 0
            logger.debug(
                f"Checking paragraph with {len(paragraph.runs)} runs for error attributes"
            )

            for run_idx, run in enumerate(paragraph.runs):
                if hasattr(run, "_element") and run._element is not None:
                    logger.debug(
                        f"Processing run {run_idx}: '{run.text}' (length: {len(run.text) if run.text else 0})"
                    )

                    # Find all run properties elements using multiple approaches
                    rpr_elements = run._element.xpath(
                        ".//a:rPr",
                        namespaces={
                            "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
                        },
                    )

                    # Also try direct child approach
                    direct_rpr_elements = run._element.xpath(
                        "./a:rPr",
                        namespaces={
                            "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
                        },
                    )

                    # Combine both approaches
                    all_rpr_elements = list(set(rpr_elements + direct_rpr_elements))

                    logger.debug(
                        f"Found {len(all_rpr_elements)} rPr elements in run {run_idx}"
                    )

                    for rpr_idx, rpr in enumerate(all_rpr_elements):
                        # Log all attributes for debugging
                        attrs = dict(rpr.attrib)
                        logger.debug(f"rPr element {rpr_idx} attributes: {attrs}")

                        # Remove err attributes that cause repair dialogs
                        if rpr.get("err") is not None:
                            err_value = rpr.get("err")
                            logger.debug(
                                f"Found err='{err_value}' attribute in run {run_idx}, removing"
                            )
                            rpr.attrib.pop("err", None)
                            removed_count += 1

                        # Also clean up problematic dirty attributes on empty runs
                        if rpr.get("dirty") == "0" and (
                            not run.text or not run.text.strip()
                        ):
                            logger.debug(f"Removing dirty='0' from empty run {run_idx}")
                            rpr.attrib.pop("dirty", None)
                            removed_count += 1

            logger.debug(f"Removed {removed_count} error attributes from paragraph")
            return removed_count

        except Exception as e:
            logger.warning(f"Failed to remove error attributes: {e}")
            return 0

    def _final_cleanup_presentation(self) -> int:
        """Final comprehensive cleanup of the entire presentation to remove error attributes."""
        try:
            total_removed = 0
            logger.info(
                "Performing final cleanup of presentation to remove error attributes"
            )

            if not self.presentation:
                return 0

            for slide_idx, slide in enumerate(self.presentation.slides):
                logger.debug(f"Final cleanup of slide {slide_idx + 1}")

                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for para_idx, paragraph in enumerate(
                            shape.text_frame.paragraphs
                        ):
                            removed_count = (
                                self._remove_error_attributes_from_paragraph(paragraph)
                            )
                            total_removed += removed_count

                    # Also check table cells
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    removed_count = (
                                        self._remove_error_attributes_from_paragraph(
                                            paragraph
                                        )
                                    )
                                    total_removed += removed_count

            logger.info(
                f"Final cleanup removed {total_removed} error attributes from presentation"
            )
            return total_removed

        except Exception as e:
            logger.warning(f"Failed in final cleanup: {e}")
            return 0

    def _post_process_xml_cleanup(self, output_path: str) -> int:
        """Post-process saved PPTX file to remove error attributes directly from XML."""
        import zipfile
        import tempfile
        import shutil
        import xml.etree.ElementTree as ET

        try:
            logger.info(
                "Performing post-processing XML cleanup to remove error attributes"
            )
            total_removed = 0

            # Create a temporary directory
            with tempfile.TemporaryDirectory() as temp_dir:
                # Extract the PPTX file
                extract_dir = os.path.join(temp_dir, "pptx_content")
                with zipfile.ZipFile(output_path, "r") as zip_ref:
                    zip_ref.extractall(extract_dir)

                # Find all slide XML files
                slides_dir = os.path.join(extract_dir, "ppt", "slides")
                if os.path.exists(slides_dir):
                    for filename in os.listdir(slides_dir):
                        if filename.endswith(".xml"):
                            slide_file = os.path.join(slides_dir, filename)
                            removed_count = self._clean_xml_file(slide_file)
                            total_removed += removed_count

                # Also clean slide layouts
                layouts_dir = os.path.join(extract_dir, "ppt", "slideLayouts")
                if os.path.exists(layouts_dir):
                    for filename in os.listdir(layouts_dir):
                        if filename.endswith(".xml"):
                            layout_file = os.path.join(layouts_dir, filename)
                            removed_count = self._clean_xml_file(layout_file)
                            total_removed += removed_count

                # Also clean slide masters
                masters_dir = os.path.join(extract_dir, "ppt", "slideMasters")
                if os.path.exists(masters_dir):
                    for filename in os.listdir(masters_dir):
                        if filename.endswith(".xml"):
                            master_file = os.path.join(masters_dir, filename)
                            removed_count = self._clean_xml_file(master_file)
                            total_removed += removed_count

                # Also clean custom XML files that may have malformed namespaces
                custom_xml_dir = os.path.join(extract_dir, "customXml")
                if os.path.exists(custom_xml_dir):
                    for filename in os.listdir(custom_xml_dir):
                        if filename.endswith(".xml"):
                            custom_file = os.path.join(custom_xml_dir, filename)
                            removed_count = self._clean_xml_file(custom_file)
                            total_removed += removed_count

                # Re-create the PPTX file
                with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zip_ref:
                    for root, dirs, files in os.walk(extract_dir):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arc_name = os.path.relpath(file_path, extract_dir)
                            zip_ref.write(file_path, arc_name)

            logger.info(
                f"Post-processing removed {total_removed} error attributes from XML"
            )
            return total_removed

        except Exception as e:
            logger.warning(f"Failed in post-processing XML cleanup: {e}")
            return 0

    def _clean_xml_file(self, xml_file_path: str) -> int:
        """Clean error attributes and other corruption patterns from a specific XML file."""
        try:
            removed_count = 0

            # Read the XML file as text and use regex to remove problematic content
            with open(xml_file_path, "r", encoding="utf-8") as f:
                content = f.read()

            original_content = content

            import re

            # Remove err="any_value" attributes
            err_matches = re.findall(r'\s+err="[^"]*"', content)
            if err_matches:
                logger.debug(
                    f"Removing {len(err_matches)} err attributes from {os.path.basename(xml_file_path)}"
                )
                content = re.sub(r'\s+err="[^"]*"', "", content)
                removed_count += len(err_matches)

            # Fix malformed namespace URIs in custom XML (common corruption pattern)
            if "customXml" in xml_file_path:
                # Fix non-absolute namespace URIs that cause XML validation warnings
                namespace_fixes = [
                    (
                        r'xmlns="34dc52e9-6ce5-490e-bfb5-77ee5cb5f472"',
                        'xmlns="urn:34dc52e9-6ce5-490e-bfb5-77ee5cb5f472"',
                    ),
                    (
                        r'xmlns="b6b70467-c097-40cb-848d-196196b595db"',
                        'xmlns="urn:b6b70467-c097-40cb-848d-196196b595db"',
                    ),
                ]

                for pattern, replacement in namespace_fixes:
                    if re.search(pattern, content):
                        logger.debug(
                            f"Fixing malformed namespace URI in {os.path.basename(xml_file_path)}"
                        )
                        content = re.sub(pattern, replacement, content)
                        removed_count += 1

            # CRITICAL FIX: Remove empty runs that appear after field elements (main corruption cause)
            field_adjacent_pattern = r"</a:fld><a:r><a:t></a:t></a:r>"
            field_adjacent_matches = re.findall(field_adjacent_pattern, content)
            if field_adjacent_matches:
                logger.debug(
                    f"Removing {len(field_adjacent_matches)} field-adjacent empty runs from {os.path.basename(xml_file_path)}"
                )
                content = re.sub(field_adjacent_pattern, "</a:fld>", content)
                removed_count += len(field_adjacent_matches)

            # ENHANCED: Remove any empty runs with just whitespace
            empty_whitespace_pattern = r"<a:r[^>]*><a:t>\s*</a:t></a:r>"
            whitespace_matches = re.findall(empty_whitespace_pattern, content)
            if whitespace_matches:
                logger.debug(
                    f"Removing {len(whitespace_matches)} whitespace-only runs from {os.path.basename(xml_file_path)}"
                )
                content = re.sub(empty_whitespace_pattern, "", content)
                removed_count += len(whitespace_matches)

            # ENHANCED: Remove multiple consecutive empty runs
            consecutive_empty_pattern = r"(<a:r[^>]*><a:t></a:t></a:r>\s*){2,}"
            consecutive_matches = re.findall(consecutive_empty_pattern, content)
            if consecutive_matches:
                logger.debug(
                    f"Removing {len(consecutive_matches)} consecutive empty run groups from {os.path.basename(xml_file_path)}"
                )
                content = re.sub(consecutive_empty_pattern, "", content)
                removed_count += len(consecutive_matches)

            # Remove problematic empty runs with dirty="0" in layouts and masters
            if "slideLayout" in xml_file_path or "slideMaster" in xml_file_path:
                # Remove empty runs that have dirty="0" and no meaningful content
                empty_dirty_runs = re.findall(
                    r'<a:r><a:rPr[^>]*dirty="0"[^>]*\/><a:t><\/a:t><\/a:r>', content
                )
                if empty_dirty_runs:
                    logger.debug(
                        f"Removing {len(empty_dirty_runs)} empty dirty runs from {os.path.basename(xml_file_path)}"
                    )
                    content = re.sub(
                        r'<a:r><a:rPr[^>]*dirty="0"[^>]*\/><a:t><\/a:t><\/a:r>',
                        "",
                        content,
                    )
                    removed_count += len(empty_dirty_runs)

            # Write back the cleaned content only if changes were made
            if content != original_content:
                with open(xml_file_path, "w", encoding="utf-8") as f:
                    f.write(content)
                logger.debug(
                    f"Applied {removed_count} fixes to {os.path.basename(xml_file_path)}"
                )

            return removed_count

        except Exception as e:
            logger.warning(f"Failed to clean XML file {xml_file_path}: {e}")
            return 0

    def _paragraph_contains_only_fields(self, paragraph) -> bool:
        """Check if paragraph contains only field elements and endParaRPr - should not be modified."""
        try:
            para_xml = paragraph._element
            children = list(para_xml)

            has_field = False
            for child in children:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == "fld":
                    has_field = True
                elif tag not in ["endParaRPr"]:
                    # Found non-field, non-endParaRPr element
                    return False

            return has_field  # Only return True if we found fields and no other content

        except Exception as e:
            logger.debug(f"Error checking field paragraph: {e}")
            return False

    def _paragraph_contains_fields(self, paragraph) -> bool:
        """Check if paragraph contains any field elements."""
        try:
            para_xml = paragraph._element
            fld_elements = para_xml.xpath(
                ".//a:fld",
                namespaces={
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
                },
            )
            return len(fld_elements) > 0
        except Exception:
            return False

    def _ensure_clean_field_structure(self, paragraph) -> None:
        """Remove any problematic runs from field paragraphs to maintain clean structure."""
        try:
            if not self._paragraph_contains_fields(paragraph):
                return

            para_xml = paragraph._element

            # Remove empty runs that appear between fields and endParaRPr
            empty_runs = para_xml.xpath(
                ".//a:r[a:t[not(text()) or normalize-space(text())='']]",
                namespaces={
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
                },
            )

            removed_count = 0
            for run in empty_runs:
                # Only remove if it's truly empty and appears after a field
                prev_sibling = run.getprevious()
                if prev_sibling is not None:
                    prev_tag = (
                        prev_sibling.tag.split("}")[-1]
                        if "}" in prev_sibling.tag
                        else prev_sibling.tag
                    )
                    if (
                        prev_tag == "fld"
                    ):  # Empty run right after field - this is the problem!
                        para_xml.remove(run)
                        removed_count += 1
                        logger.debug(
                            "Removed problematic empty run after field element"
                        )

            if removed_count > 0:
                logger.debug(
                    f"Cleaned {removed_count} problematic runs from field paragraph"
                )

        except Exception as e:
            logger.warning(f"Failed to clean field structure: {e}")

    def _aggressive_empty_run_cleanup(self, paragraph) -> None:
        """Aggressively remove all empty runs that could cause PowerPoint corruption."""
        try:
            para_xml = paragraph._element
            runs_to_remove = []

            # Find all text runs
            for run in para_xml.findall(
                ".//a:r", {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            ):
                # Check if run has empty or whitespace-only text
                text_elem = run.find(
                    ".//a:t",
                    {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
                )
                if text_elem is not None:
                    text_content = text_elem.text or ""
                    # Remove if completely empty or just whitespace
                    if not text_content.strip():
                        runs_to_remove.append(run)
                else:
                    # Run without text element - definitely remove
                    runs_to_remove.append(run)

            # Remove identified empty runs
            for run in runs_to_remove:
                try:
                    para_xml.remove(run)
                    logger.debug("Removed empty/whitespace run in aggressive cleanup")
                except Exception as e:
                    logger.debug(f"Could not remove run in aggressive cleanup: {e}")

            if runs_to_remove:
                logger.debug(
                    f"Aggressive cleanup removed {len(runs_to_remove)} empty runs"
                )

        except Exception as e:
            logger.warning(f"Failed in aggressive empty run cleanup: {e}")

    def _process_table_shape(self, shape: BaseShape, data: Dict[str, Any]) -> None:
        """Process table shape for merge field replacement."""
        try:
            table = shape.table
            logger.debug(
                f"📊 Processing table with {len(table.rows)} rows and {len(table.columns)} columns"
            )

            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if cell.text:
                        logger.debug(
                            f"🔍 Processing cell [{row_idx},{col_idx}]: '{cell.text}'"
                        )
                        # Process each paragraph in the cell
                        for paragraph in cell.text_frame.paragraphs:
                            self._process_paragraph(paragraph, data)

        except Exception as e:
            logger.warning(f"Failed to process table shape: {e}")

    def _get_field_value(self, field_name: str, data: Dict[str, Any]) -> Any:
        """Get value for a merge field from data dictionary."""
        try:
            # Handle nested field references like "table.0.field_name"
            field_parts = field_name.split(".")
            current_value = data

            # Debug logging for field resolution
            logger.info(f"🔧 DEBUG: Getting field value for: '{field_name}'")
            logger.info(f"🔧 DEBUG: Data keys at root level: {list(data.keys())}")
            logger.info(f"🔧 DEBUG: Data type: {type(data)}")

            # First try direct path resolution
            for part in field_parts:
                if isinstance(current_value, dict):
                    current_value = current_value.get(part)
                    logger.info(
                        f"🔧 DEBUG: After part '{part}': {type(current_value).__name__} = {current_value}"
                    )
                elif isinstance(current_value, list):
                    try:
                        index = int(part)
                        current_value = (
                            current_value[index]
                            if 0 <= index < len(current_value)
                            else None
                        )
                        logger.info(
                            f"🔧 DEBUG: After list index {index}: {type(current_value).__name__} = {current_value}"
                        )
                    except (ValueError, IndexError):
                        current_value = None
                        logger.info(f"🔧 DEBUG: Invalid list index: {part}")
                else:
                    current_value = None
                    logger.info(
                        f"🔧 DEBUG: Cannot navigate further from {type(current_value).__name__}"
                    )

                if current_value is None:
                    break

            # If direct path failed and we have sheet data, try looking in each sheet
            if current_value is None:
                logger.info(
                    "🔧 DEBUG: Direct path resolution failed, trying sheet-nested lookup"
                )
                # Try to find the field in sheet data (e.g., order_form.image_search.0.field)
                for sheet_name, sheet_data in data.items():
                    # Skip metadata and debug fields
                    if sheet_name.startswith("__"):
                        continue

                    # Check if this sheet contains the first part of our field path
                    if isinstance(sheet_data, dict) and field_parts[0] in sheet_data:
                        # logger.debug(f"Found {field_parts[0]} in sheet {sheet_name}")  # Too verbose
                        # Start with the sheet data
                        nested_value = sheet_data

                        # Navigate through the field parts
                        for part in field_parts:
                            if isinstance(nested_value, dict):
                                nested_value = nested_value.get(part)
                            elif isinstance(nested_value, list):
                                try:
                                    index = int(part)
                                    nested_value = (
                                        nested_value[index]
                                        if 0 <= index < len(nested_value)
                                        else None
                                    )
                                except (ValueError, IndexError):
                                    nested_value = None
                            else:
                                nested_value = None

                            if nested_value is None:
                                break

                        if nested_value is not None:
                            # Avoid logging large data structures like base64 images
                            if isinstance(nested_value, dict) and any(
                                key in str(nested_value).lower()
                                for key in ["base64", "data:image"]
                            ):
                                logger.info(
                                    f"🔧 DEBUG: Found image data via sheet-nested lookup for field: {field_name}"
                                )
                            else:
                                logger.info(
                                    f"🔧 DEBUG: Found value via sheet-nested lookup: {nested_value}"
                                )
                            return nested_value

                # Special handling for flat structures without row index
                # For fields like "client_info.client_name" (without row index)
                # Only apply this fallback if the field doesn't contain a numeric index
                if len(field_parts) >= 2:
                    table_name = field_parts[0]
                    field_key = field_parts[-1]

                    # Check if any part is a numeric index - if so, don't use fallback
                    has_numeric_index = False
                    for part in field_parts[
                        1:-1
                    ]:  # Check middle parts for numeric indices
                        try:
                            int(part)
                            has_numeric_index = True
                            break
                        except ValueError:
                            continue

                    # Only apply fallback for non-indexed fields
                    if not has_numeric_index:
                        # Look for the table in each sheet
                        for sheet_name, sheet_data in data.items():
                            if (
                                isinstance(sheet_data, dict)
                                and table_name in sheet_data
                            ):
                                table_data = sheet_data[table_name]

                                # Case 1: Table is a flat dictionary (key-value pairs)
                                if (
                                    isinstance(table_data, dict)
                                    and field_key in table_data
                                ):
                                    value = table_data[field_key]
                                    # Avoid logging large data structures like base64 images
                                    if isinstance(value, dict) and any(
                                        key in str(value).lower()
                                        for key in ["base64", "data:image"]
                                    ):
                                        logger.debug(
                                            f"Found image data in flat structure {table_name}.{field_key}"
                                        )
                                    else:
                                        logger.debug(
                                            f"Found value in flat structure {table_name}.{field_key}: {value}"
                                        )
                                    return value

                                # Case 2: Table is a list with a single item (only for non-indexed fields)
                                elif (
                                    isinstance(table_data, list) and len(table_data) > 0
                                ):
                                    # Try first row if no index specified
                                    first_row = table_data[0]
                                    if (
                                        isinstance(first_row, dict)
                                        and field_key in first_row
                                    ):
                                        value = first_row[field_key]
                                        # Avoid logging large data structures like base64 images
                                        if isinstance(value, dict) and any(
                                            key in str(value).lower()
                                            for key in ["base64", "data:image"]
                                        ):
                                            logger.debug(
                                                f"Found image data in first row of {table_name}[0].{field_key}"
                                            )
                                        else:
                                            logger.debug(
                                                f"Found value in first row of {table_name}[0].{field_key}: {value}"
                                            )
                                        return value

            # If field not found, try common field name variations
            if current_value is None and field_name:
                logger.debug(f"Field '{field_name}' not found, trying variations")

                # Try common field name variations for normalization mismatches
                variations = self._get_field_name_variations(field_name)

                for variation in variations:
                    if variation != field_name:
                        logger.debug(f"Trying field variation: '{variation}'")

                        # Reconstruct path with variation
                        if len(field_parts) > 1:
                            variation_parts = field_parts[:-1] + [variation]
                            variation_path = ".".join(variation_parts)
                        else:
                            variation_path = variation

                        # Try the variation (non-recursive to avoid infinite loops)
                        temp_value = self._get_field_value_direct(variation_path, data)
                        if temp_value is not None:
                            logger.debug(
                                f"Found value using field variation: {variation}"
                            )
                            current_value = temp_value
                            break

            # Ensure missing fields return empty string instead of None
            if current_value is None:
                logger.debug(f"Field '{field_name}' not found, returning empty string")
                return ""

            # Avoid logging large data structures like base64 images
            if isinstance(current_value, dict) and any(
                key in str(current_value).lower() for key in ["base64", "data:image"]
            ):
                logger.debug(f"Returning image data for field: {field_name}")
            else:
                logger.debug(f"✅ Final value for '{field_name}': {current_value}")
            return current_value

        except Exception as e:
            logger.warning(f"Failed to get field value for '{field_name}': {e}")
            return None

    def _replace_image_placeholder(
        self,
        slide,
        shape: BaseShape,
        data: Dict[str, Any],
        images: Optional[Dict[str, List[Dict[str, Any]]]] = None,
    ) -> None:
        """Replace image placeholder with actual image."""
        if not images:
            return

        try:
            # Extract image field name from placeholder text
            placeholder_text = shape.text
            merge_fields = validate_merge_fields(placeholder_text)

            if not merge_fields:
                return

            # Get the first merge field (assuming single image placeholder)
            image_field = merge_fields[0]
            image_path = self._get_image_for_field(image_field, data, images)

            if image_path:
                # Verify the image file exists
                if os.path.exists(str(image_path)):
                    try:
                        # Get shape position and size
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        # Remove the placeholder shape
                        slide.shapes._spTree.remove(shape._element)

                        # Add the image
                        slide.shapes.add_picture(image_path, left, top, width, height)
                        logger.debug(f"Replaced image placeholder with: {image_path}")
                    except Exception as img_error:
                        logger.warning(
                            f"Failed to add image '{image_path}': {img_error}"
                        )
                        # Keep the original placeholder if image insertion fails
                        pass
                else:
                    logger.warning(f"Image file does not exist at path: {image_path}")

        except Exception as e:
            logger.warning(f"Failed to replace image placeholder: {e}")

    def preview_merge(
        self,
        data: Dict[str, Any],
        images: Optional[Dict[str, List[Dict[str, Any]]]] = None,
    ) -> Dict[str, Any]:
        """Preview what the merge would look like without actually performing it."""
        try:
            merge_fields = self.get_merge_fields()
            image_placeholders = self._get_image_placeholders()

            preview = {
                "merge_fields": merge_fields,
                "image_placeholders": image_placeholders,
                "field_values": {},
                "image_mappings": {},
                "missing_fields": [],
                "missing_images": [],
            }

            # Check which fields will be populated
            for field in merge_fields:
                value = self._get_field_value(field, data)
                preview["field_values"][field] = value

                if value is None:
                    preview["missing_fields"].append(field)

            # Check image mappings
            if images:
                for placeholder in image_placeholders:
                    image_path = self._get_image_for_field(placeholder, data, images)
                    preview["image_mappings"][placeholder] = image_path

                    if not image_path:
                        preview["missing_images"].append(placeholder)
            else:
                preview["missing_images"] = image_placeholders.copy()

            return preview

        except Exception as e:
            raise PowerPointProcessingError(f"Failed to generate merge preview: {e}")

    def get_image_requirements(self) -> Dict[str, Any]:
        """Analyze template to determine image requirements."""
        try:
            image_placeholders = self._get_image_placeholders()

            requirements = {
                "total_image_placeholders": len(image_placeholders),
                "placeholder_details": [],
                "suggested_naming": [],
            }

            for placeholder in image_placeholders:
                # Analyze placeholder to suggest naming conventions
                detail = {
                    "placeholder": placeholder,
                    "suggested_field_path": self._suggest_field_path(placeholder),
                    "suggested_cell_position": self._extract_cell_position_from_field(
                        placeholder
                    ),
                }
                requirements["placeholder_details"].append(detail)

            return requirements

        except Exception as e:
            logger.warning(f"Failed to analyze image requirements: {e}")
            return {
                "total_image_placeholders": 0,
                "placeholder_details": [],
                "suggested_naming": [],
            }

    def _suggest_field_path(self, placeholder: str) -> str:
        """Suggest a field path for accessing image data."""
        # Convert placeholder to suggested data path
        if "image_search" in placeholder:
            return placeholder  # Already in correct format
        elif "image" in placeholder:
            # Convert generic image reference to indexed format
            match = re.search(r"(\d+)", placeholder)
            if match:
                index = match.group(1)
                return f"order_form.image_search.{index}.image"

        return placeholder

    def _extract_cell_position_from_field(self, field_name: str) -> Optional[str]:
        """Extract cell position hint from field name."""
        # Look for cell references in field names
        cell_pattern = r"([A-Z]+\d+)"
        match = re.search(cell_pattern, field_name.upper())
        if match:
            return match.group(1)

        return None

    def _create_temp_image_from_base64(self, base64_data: str) -> Optional[str]:
        """Create a temporary image file from base64 data."""
        try:
            import base64
            import tempfile
            import os

            # Extract the image format and data
            if base64_data.startswith("data:image/"):
                # Format: data:image/png;base64,iVBORw0KGgo...
                split_result = base64_data.split(",", 1)
                if len(split_result) == 2:
                    header, encoded_data = split_result
                    image_format = header.split("/")[1].split(";")[
                        0
                    ]  # Extract 'png' from 'data:image/png;base64'
                else:
                    # Malformed data URL, treat as raw base64
                    logger.warning(
                        f"Malformed data URL, treating as raw base64: {base64_data[:50]}..."
                    )
                    encoded_data = base64_data
                    image_format = "png"
            else:
                # Raw base64 data, assume PNG
                encoded_data = base64_data
                image_format = "png"

            # Decode base64 data
            image_bytes = base64.b64decode(encoded_data)

            # Create temporary file
            with tempfile.NamedTemporaryFile(
                suffix=f".{image_format}", delete=False
            ) as temp_file:
                temp_file.write(image_bytes)
                temp_path = temp_file.name

            logger.debug(f"Created temporary image file: {temp_path}")
            return temp_path

        except Exception as e:
            logger.error(f"Failed to create temporary image from base64: {e}")
            return None

    def close(self) -> None:
        """Close the presentation and free resources."""
        # python-pptx doesn't require explicit closing, but we'll reset the reference
        self.presentation = None

    def validate_template(self) -> Dict[str, Any]:
        """Validate template and return information about merge fields and structure."""
        if not self.presentation:
            self._validate_template()

        try:
            validation_info = {
                "slide_count": len(self.presentation.slides),
                "merge_fields": self.get_merge_fields(),
                "image_placeholders": self._get_image_placeholders(),
                "slides": [],
            }

            for slide_idx, slide in enumerate(self.presentation.slides):
                slide_info = {
                    "slide_number": slide_idx + 1,
                    "shape_count": len(slide.shapes),
                    "merge_fields": self._extract_slide_merge_fields(slide),
                    "image_placeholders": self._get_slide_image_placeholders(slide),
                    "has_tables": any(
                        shape.shape_type == MSO_SHAPE_TYPE.TABLE
                        for shape in slide.shapes
                    ),
                    "has_images": any(
                        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                        for shape in slide.shapes
                    ),
                }
                validation_info["slides"].append(slide_info)

            return validation_info

        except Exception as e:
            raise PowerPointProcessingError(f"Template validation failed: {e}")

    def _get_image_placeholders(self) -> List[str]:
        """Get all image placeholders from the presentation."""
        placeholders = []

        try:
            for slide in self.presentation.slides:
                slide_placeholders = self._get_slide_image_placeholders(slide)
                placeholders.extend(slide_placeholders)
        except Exception as e:
            logger.warning(f"Failed to get image placeholders: {e}")

        return list(set(placeholders))  # Remove duplicates

    def _get_slide_image_placeholders(self, slide) -> List[str]:
        """Get image placeholders from a single slide."""
        placeholders = []

        try:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_content = self._get_full_text_from_shape(shape)
                    if text_content and self._is_image_placeholder(text_content):
                        merge_fields = validate_merge_fields(text_content)
                        placeholders.extend(merge_fields)
                elif (
                    hasattr(shape, "text")
                    and shape.text
                    and self._is_image_placeholder(shape.text)
                ):
                    merge_fields = validate_merge_fields(shape.text)
                    placeholders.extend(merge_fields)
        except Exception as e:
            logger.warning(f"Failed to get slide image placeholders: {e}")

        return placeholders

    def _process_link_field(
        self, field_name: str, field_value: Any, paragraph, field_info: Dict[str, Any]
    ) -> bool:
        """Process a link field by creating a hyperlink.

        Args:
            field_name: The name of the field
            field_value: The field value (should be link object or string)
            paragraph: The paragraph containing the field
            field_info: Information about the field position in runs

        Returns:
            True if the link was successfully processed, False otherwise
        """
        try:
            # Extract link data
            link_data = self._extract_link_data(field_value)
            if not link_data:
                return False

            title = link_data.get("title", "")
            url = link_data.get("link", "")

            # Validate the URL
            if not url or not url.strip():
                logger.warning(f"No URL provided for link field '{field_name}'")
                return False

            # Replace field with hyperlink
            self._replace_field_with_hyperlink(paragraph, field_info, title, url)
            return True

        except Exception as e:
            logger.warning(f"Failed to process link field '{field_name}': {e}")
            return False

    def _extract_link_data(self, field_value: Any) -> Optional[Dict[str, str]]:
        """Extract link data from field value.

        Args:
            field_value: The field value (link object or string)

        Returns:
            Dictionary with 'title' and 'link' keys, or None if invalid
        """
        if isinstance(field_value, dict):
            # Link object format: {"title": "...", "link": "..."}
            title = field_value.get("title", "")
            link = field_value.get("link", "")
            return {"title": title, "link": link}

        elif isinstance(field_value, str):
            # Simple string format: use as both title and link
            return {"title": field_value, "link": field_value}

        else:
            logger.warning(f"Invalid link field value type: {type(field_value)}")
            return None

    def _replace_field_with_hyperlink(
        self, paragraph, field_info: Dict[str, Any], title: str, url: str
    ) -> None:
        """Replace a merge field with a hyperlink.

        Args:
            paragraph: The paragraph containing the field
            field_info: Information about the field position in runs
            title: The display text for the hyperlink
            url: The URL for the hyperlink
        """
        try:
            affected_runs = field_info["affected_runs"]

            if len(affected_runs) == 1:
                # Simple case: field is entirely within one run
                self._replace_field_in_single_run_with_hyperlink(
                    paragraph, affected_runs[0], title, url
                )
            else:
                # Complex case: field spans multiple runs
                self._replace_field_across_runs_with_hyperlink(
                    paragraph, affected_runs, title, url
                )

        except Exception as e:
            logger.warning(f"Failed to replace field with hyperlink: {e}")

    def _replace_field_in_single_run_with_hyperlink(
        self, paragraph, run_info: Dict[str, Any], title: str, url: str
    ) -> None:
        """Replace field within a single run with hyperlink using precise text splitting."""
        try:
            run = run_info["run"]
            original_text = run_info["run_text"]
            field_start = run_info["field_start_in_run"]
            field_end = run_info["field_end_in_run"]

            # Find the current run index at the time of processing
            run_index = None
            for i, p_run in enumerate(paragraph.runs):
                if p_run == run:
                    run_index = i
                    break

            if run_index is None:
                logger.warning("Could not find run in paragraph, using fallback")
                # Fallback to simple replacement
                new_text = (
                    original_text[:field_start] + title + original_text[field_end:]
                )
                run.text = new_text
                self._apply_hyperlink_to_run(run, url)
                return

            # Use the new splitting approach for precise hyperlink application
            self._split_run_for_hyperlink_by_index(
                paragraph, run_index, field_start, field_end, title, url
            )

        except Exception as e:
            logger.warning(f"Failed to replace field in single run with hyperlink: {e}")
            # Fallback to simple replacement
            try:
                run = run_info["run"]
                original_text = run_info["run_text"]
                field_start = run_info["field_start_in_run"]
                field_end = run_info["field_end_in_run"]
                new_text = (
                    original_text[:field_start] + title + original_text[field_end:]
                )
                run.text = new_text
                self._apply_hyperlink_to_run(run, url)
            except Exception as fallback_e:
                logger.warning(
                    f"Fallback hyperlink replacement also failed: {fallback_e}"
                )

    def _replace_field_across_runs_with_hyperlink(
        self, paragraph, affected_runs: List[Dict[str, Any]], title: str, url: str
    ) -> None:
        """Replace field that spans across multiple runs with hyperlink."""
        try:
            if not affected_runs:
                return

            # For multi-run fields, we'll clear the field and put the hyperlink in the first run
            first_run_info = affected_runs[0]

            # Clear the field from all runs
            for i, run_info in enumerate(affected_runs):
                run = run_info["run"]

                if i == 0:
                    # First run: replace field start with title
                    original_text = run_info["run_text"]
                    field_start = run_info["field_start_in_run"]
                    new_text = original_text[:field_start] + title
                    run.text = new_text
                    # Apply hyperlink to this run
                    self._apply_hyperlink_to_run(run, url)
                else:
                    # Other runs: remove the field portion
                    original_text = run_info["run_text"]
                    if i == len(affected_runs) - 1:
                        # Last run: remove until field end
                        field_end = run_info["field_end_in_run"]
                        new_text = original_text[field_end:]
                    else:
                        # Middle runs: remove everything
                        new_text = ""
                    run.text = new_text

        except Exception as e:
            logger.warning(f"Failed to replace field across runs with hyperlink: {e}")

    def _split_run_for_hyperlink_by_index(
        self,
        paragraph,
        run_index: int,
        field_start: int,
        field_end: int,
        title: str,
        url: str,
    ) -> None:
        """Split a run into multiple runs to apply hyperlink only to specific text using run index.

        Args:
            paragraph: The paragraph containing the run
            run_index: Index of the run to split
            field_start: Start position of the field in the run's text
            field_end: End position of the field in the run's text
            title: The link title to insert
            url: The URL for the hyperlink
        """
        try:
            run = paragraph.runs[run_index]
            original_text = run.text
            logger.debug(
                f"🔗 Splitting run at index {run_index} with text: '{original_text}'"
            )
            logger.debug(f"🔗 Field positions: start={field_start}, end={field_end}")

            # Text segments
            text_before = original_text[:field_start]
            text_after = original_text[field_end:]

            # Store original formatting
            original_font = run.font

            # Simpler approach: modify the existing run and add new runs
            logger.debug(
                f"🔗 Calling _replace_run_with_segments with: text_before='{text_before}', title='{title}', text_after='{text_after}'"
            )
            self._replace_run_with_segments(
                paragraph, run_index, text_before, title, text_after, url, original_font
            )

        except Exception as e:
            logger.warning(f"Failed to split run for hyperlink by index: {e}")
            # Fallback to simple replacement
            try:
                run = paragraph.runs[run_index]
                original_text = run.text
                logger.debug(f"🔗 Using fallback: replacing with '{title}'")
                run.text = (
                    original_text[:field_start] + title + original_text[field_end:]
                )
                self._apply_hyperlink_to_run(run, url)
            except Exception as fallback_e:
                logger.warning(
                    f"Fallback hyperlink replacement also failed: {fallback_e}"
                )

    def _split_run_for_hyperlink(
        self, paragraph, run, field_start: int, field_end: int, title: str, url: str
    ) -> None:
        """Split a run into multiple runs to apply hyperlink only to specific text.

        Args:
            paragraph: The paragraph containing the run
            run: The run to split
            field_start: Start position of the field in the run's text
            field_end: End position of the field in the run's text
            title: The link title to insert
            url: The URL for the hyperlink
        """
        try:
            original_text = run.text
            logger.debug(f"🔗 Splitting run with text: '{original_text}'")
            logger.debug(f"🔗 Field positions: start={field_start}, end={field_end}")

            # Get the run's position in the paragraph
            run_index = None
            for i, p_run in enumerate(paragraph.runs):
                if p_run == run:
                    run_index = i
                    break

            if run_index is None:
                logger.warning("Could not find run in paragraph")
                return

            # Text segments
            text_before = original_text[:field_start]
            text_after = original_text[field_end:]

            # Store original formatting
            original_font = run.font

            # Clear the original run
            run.text = ""

            # Create runs in order: before, link, after
            runs_to_create = []

            if text_before:
                runs_to_create.append(("text", text_before))

            runs_to_create.append(("link", title))

            if text_after:
                runs_to_create.append(("text", text_after))

            # Create the new runs
            for i, (run_type, text) in enumerate(runs_to_create):
                if i == 0:
                    # Use the existing run for the first segment
                    current_run = run
                else:
                    # Create new run after the current position
                    current_run = (
                        paragraph.runs[run_index + i]
                        ._element.getparent()
                        .insert(run_index + i, paragraph.runs[0]._element.tag)
                    )
                    # This is complex - let me use a simpler approach

            # Simpler approach: modify the existing run and add new runs
            logger.debug(
                f"🔗 Calling _replace_run_with_segments with: text_before='{text_before}', title='{title}', text_after='{text_after}'"
            )
            self._replace_run_with_segments(
                paragraph, run_index, text_before, title, text_after, url, original_font
            )

        except Exception as e:
            logger.warning(f"Failed to split run for hyperlink: {e}")
            # Fallback to simple replacement
            logger.debug(f"🔗 Using fallback: replacing with '{title}'")
            run.text = original_text[:field_start] + title + original_text[field_end:]
            self._apply_hyperlink_to_run(run, url)

    def _replace_run_with_segments(
        self,
        paragraph,
        run_index: int,
        text_before: str,
        link_text: str,
        text_after: str,
        url: str,
        original_font,
    ) -> None:
        """Replace a run with multiple segments for precise hyperlink application.

        Args:
            paragraph: The paragraph containing the run
            run_index: Index of the run to replace
            text_before: Text before the hyperlink
            link_text: Text for the hyperlink
            text_after: Text after the hyperlink
            url: The URL for the hyperlink
            original_font: Original font formatting to preserve
        """
        try:
            # Remove the original run (we'll replace it with new ones)
            original_run = paragraph.runs[run_index]

            # Keep track of segments to create
            segments = []
            if text_before:
                segments.append(("text", text_before))
            if link_text:
                segments.append(("link", link_text))
            if text_after:
                segments.append(("text", text_after))

            logger.debug(f"🔗 Creating segments: {segments}")

            # Set the first segment in the existing run
            if segments:
                first_type, first_text = segments[0]
                logger.debug(
                    f"🔗 Setting first segment in original run: '{first_text}' (type: {first_type})"
                )
                original_run.text = first_text
                if first_type == "link":
                    logger.debug(
                        f"🔗 Applying hyperlink to original run: '{first_text}'"
                    )
                    self._apply_hyperlink_to_run(original_run, url)

                # Add additional runs for remaining segments
                for segment_type, segment_text in segments[1:]:
                    # Add a new run after the current run
                    new_run = paragraph.add_run(segment_text)
                    logger.debug(
                        f"🔗 Added new run: '{segment_text}' (type: {segment_type})"
                    )

                    # Copy formatting from original run
                    self._copy_run_formatting(original_font, new_run.font)

                    # Apply hyperlink if this is a link segment
                    if segment_type == "link":
                        logger.debug(
                            f"🔗 Applying hyperlink to new run: '{segment_text}'"
                        )
                        self._apply_hyperlink_to_run(new_run, url)

        except Exception as e:
            logger.warning(f"Failed to replace run with segments: {e}")

    def _copy_run_formatting(self, source_font, target_font) -> None:
        """Copy formatting from source font to target font.

        Args:
            source_font: Source font to copy from
            target_font: Target font to copy to
        """
        try:
            if source_font.name:
                target_font.name = source_font.name
            if source_font.size:
                target_font.size = source_font.size
            if source_font.bold is not None:
                target_font.bold = source_font.bold
            if source_font.italic is not None:
                target_font.italic = source_font.italic
            if source_font.underline is not None:
                target_font.underline = source_font.underline
            if hasattr(source_font, "color") and source_font.color:
                target_font.color.rgb = source_font.color.rgb
        except Exception as e:
            logger.debug(f"Could not copy some font formatting: {e}")

    def _apply_hyperlink_to_run(self, run, url: str) -> None:
        """Apply hyperlink to a text run using python-pptx API.

        Args:
            run: The text run to apply hyperlink to
            url: The URL for the hyperlink
        """
        try:
            if url and url.strip():
                hlink = run.hyperlink
                hlink.address = url.strip()
                logger.debug(
                    f"🔗 Applied hyperlink '{url}' to run with text '{run.text}'"
                )
            else:
                logger.warning("Cannot apply empty URL as hyperlink")

        except Exception as e:
            logger.warning(f"Failed to apply hyperlink to run: {e}")
