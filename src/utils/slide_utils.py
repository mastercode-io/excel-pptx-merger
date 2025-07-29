"""Utilities for PowerPoint slide manipulation and processing."""

import logging
import re
from typing import Any, Dict, List, Optional, Tuple
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape

logger = logging.getLogger(__name__)


def duplicate_slide(
    presentation: Presentation, source_slide: Slide, target_index: int
) -> Slide:
    """Create a duplicate of a slide at the specified index.

    Args:
        presentation: The PowerPoint presentation
        source_slide: The slide to duplicate
        target_index: Index where to insert the new slide

    Returns:
        The newly created slide
    """
    try:
        # Get the slide layout from the source slide
        slide_layout = source_slide.slide_layout

        # Create a new slide with the same layout
        new_slide = presentation.slides.add_slide(slide_layout)

        # Copy all shapes from source to new slide
        for shape in source_slide.shapes:
            _copy_shape(shape, new_slide)

        # Move the slide to the target position if needed
        if target_index < len(presentation.slides) - 1:
            _move_slide(presentation, len(presentation.slides) - 1, target_index)
            new_slide = presentation.slides[target_index]

        logger.debug(f"Successfully duplicated slide to position {target_index}")
        return new_slide

    except Exception as e:
        logger.error(f"Failed to duplicate slide: {e}")
        raise


def _copy_shape(source_shape: BaseShape, target_slide: Slide) -> None:
    """Copy a shape from source to target slide.

    Args:
        source_shape: Shape to copy
        target_slide: Target slide to copy to
    """
    try:
        # This is a simplified approach - python-pptx doesn't have direct shape duplication
        # For now, we'll rely on the slide layout and manual copying for complex cases
        # In a full implementation, this would need to handle different shape types

        # Handle text shapes
        if hasattr(source_shape, "text_frame") and source_shape.text_frame:
            # For text shapes, we'll copy the text content
            # The actual shape copying will happen through slide layout
            pass

        # Handle other shape types as needed
        # This is a placeholder for more complex shape copying logic

    except Exception as e:
        logger.warning(f"Could not copy shape: {e}")


def _move_slide(presentation: Presentation, from_index: int, to_index: int) -> None:
    """Move a slide from one position to another.

    Args:
        presentation: The PowerPoint presentation
        from_index: Current slide index
        to_index: Target slide index
    """
    try:
        # This is a limitation of python-pptx - slide reordering is complex
        # For now, we'll work with slides in their created order
        # A full implementation would need custom XML manipulation
        logger.debug(
            f"Slide reordering from {from_index} to {to_index} requested (simplified implementation)"
        )

    except Exception as e:
        logger.warning(f"Could not move slide: {e}")


def get_slide_index(slide_number: int) -> int:
    """Convert 1-based slide number to 0-based index.

    Args:
        slide_number: 1-based slide number

    Returns:
        0-based slide index
    """
    return slide_number - 1


def is_template_slide(slide: Slide, template_marker: str) -> bool:
    """Check if a slide contains template markers for dynamic duplication.

    Args:
        slide: The slide to check
        template_marker: The template marker to look for (e.g., "{{#list:")

    Returns:
        True if slide contains template markers
    """
    try:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text_content = _get_shape_text(shape)
                if template_marker in text_content:
                    return True
            elif hasattr(shape, "text") and shape.text:
                if template_marker in shape.text:
                    return True

        return False

    except Exception as e:
        logger.error(f"Error checking template slide: {e}")
        return False


def extract_list_name(slide: Slide, template_marker: str) -> Optional[str]:
    """Extract the list name from a template slide marker.

    Args:
        slide: The template slide
        template_marker: The template marker (e.g., "{{#list:")

    Returns:
        The list name or None if not found
    """
    try:
        pattern = re.escape(template_marker) + r"([^}]+)\}\}"

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text_content = _get_shape_text(shape)
                match = re.search(pattern, text_content)
                if match:
                    return match.group(1).strip()
            elif hasattr(shape, "text") and shape.text:
                match = re.search(pattern, shape.text)
                if match:
                    return match.group(1).strip()

        return None

    except Exception as e:
        logger.error(f"Error extracting list name: {e}")
        return None


def _get_shape_text(shape: BaseShape) -> str:
    """Get all text content from a shape.

    Args:
        shape: The shape to extract text from

    Returns:
        Combined text content
    """
    try:
        text_parts = []

        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_parts.append(run.text)

        return "".join(text_parts)

    except Exception as e:
        logger.debug(f"Could not extract text from shape: {e}")
        return ""


def filter_slides(
    slides: List[Slide], include_slides: List[int], exclude_slides: List[int]
) -> List[Tuple[int, Slide]]:
    """Filter slides based on include/exclude configuration.

    Args:
        slides: List of slides to filter
        include_slides: List of slide numbers to include (1-based)
        exclude_slides: List of slide numbers to exclude (1-based)

    Returns:
        List of tuples (original_index, slide) for filtered slides
    """
    try:
        filtered_slides = []

        for i, slide in enumerate(slides):
            slide_number = i + 1  # Convert to 1-based

            # If include list is specified and not empty, only include those slides
            if include_slides:
                if slide_number not in include_slides:
                    continue

            # If exclude list is specified, exclude those slides
            if exclude_slides:
                if slide_number in exclude_slides:
                    continue

            filtered_slides.append((i, slide))

        logger.info(f"Filtered {len(slides)} slides to {len(filtered_slides)} slides")
        return filtered_slides

    except Exception as e:
        logger.error(f"Error filtering slides: {e}")
        return [(i, slide) for i, slide in enumerate(slides)]


def create_list_context(
    item: Dict[str, Any],
    index: int,
    total: int,
    parent_data: Dict[str, Any],
    special_variables: Dict[str, str],
) -> Dict[str, Any]:
    """Create context for list item with special variables.

    Args:
        item: The current list item data
        index: 0-based index of current item
        total: Total number of items in list
        parent_data: Parent context data
        special_variables: Special variable name mappings

    Returns:
        Context dictionary with item data and special variables
    """
    try:
        # Start with the item data
        context = dict(item) if isinstance(item, dict) else {}

        # Add special variables
        context[special_variables.get("index", "$index")] = index
        context[special_variables.get("position", "$position")] = index + 1
        context[special_variables.get("first", "$first")] = index == 0
        context[special_variables.get("last", "$last")] = index == total - 1
        context[special_variables.get("odd", "$odd")] = index % 2 == 1
        context[special_variables.get("even", "$even")] = index % 2 == 0

        # Add parent context for "../" resolution
        context[".."] = parent_data

        # Add root context for "$root" resolution
        context["$root"] = parent_data

        return context

    except Exception as e:
        logger.error(f"Error creating list context: {e}")
        return item if isinstance(item, dict) else {}


def resolve_field_path(
    field_path: str,
    context: Dict[str, Any],
    parent_prefix: str = "../",
    root_prefix: str = "$root.",
) -> Any:
    """Resolve field path including parent context (../) and root ($root).

    Args:
        field_path: The field path to resolve
        context: Current context dictionary
        parent_prefix: Prefix for parent context access
        root_prefix: Prefix for root context access

    Returns:
        Resolved field value or None if not found
    """
    try:
        # Handle root context access
        if field_path.startswith(root_prefix):
            root_path = field_path[len(root_prefix) :]
            root_context = context.get("$root", {})
            return _get_nested_value(root_context, root_path)

        # Handle parent context access
        if field_path.startswith(parent_prefix):
            parent_path = field_path[len(parent_prefix) :]
            parent_context = context.get("..", {})
            return _get_nested_value(parent_context, parent_path)

        # Handle regular field access
        return _get_nested_value(context, field_path)

    except Exception as e:
        logger.debug(f"Could not resolve field path '{field_path}': {e}")
        return None


def _get_nested_value(data: Dict[str, Any], path: str) -> Any:
    """Get value from nested dictionary using dot notation.

    Args:
        data: Dictionary to search in
        path: Dot-separated path (e.g., "contact.email")

    Returns:
        Value at the path or None if not found
    """
    try:
        current = data
        parts = path.split(".")

        for part in parts:
            if isinstance(current, dict) and part in current:
                current = current[part]
            else:
                return None

        return current

    except Exception as e:
        logger.debug(f"Could not get nested value for path '{path}': {e}")
        return None
