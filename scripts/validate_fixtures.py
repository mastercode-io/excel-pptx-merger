#!/usr/bin/env python3
"""
Script to validate the generated test fixtures
"""

import os
import json
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
import re


def validate_excel_file(excel_path="tests/fixtures/sample_excel.xlsx"):
    """Validate the sample Excel file structure"""
    print(f"🔍 Validating Excel file: {excel_path}")

    if not os.path.exists(excel_path):
        print(f"❌ Excel file not found: {excel_path}")
        return False

    try:
        # Test with pandas
        df = pd.read_excel(excel_path, sheet_name='Order Form', header=None)
        print(f"✅ Excel file loaded successfully")
        print(f"   📊 Shape: {df.shape}")

        # Test with openpyxl for images
        wb = load_workbook(excel_path)
        ws = wb['Order Form']

        # Check for images
        if hasattr(ws, '_images') and ws._images:
            print(f"   🖼️  Found {len(ws._images)} embedded images")
        else:
            print(f"   ⚠️  No embedded images found")

        # Display first few rows for verification
        print(f"   📝 Sample data:")
        for i in range(min(5, len(df))):
            row_data = [str(x) if pd.notna(x) else "" for x in df.iloc[i][:6]]
            print(f"      Row {i + 1}: {row_data}")

        wb.close()
        return True

    except Exception as e:
        print(f"❌ Error validating Excel file: {e}")
        return False


def validate_powerpoint_file(pptx_path="tests/fixtures/sample_template.pptx"):
    """Validate the sample PowerPoint template"""
    print(f"\n🔍 Validating PowerPoint template: {pptx_path}")

    if not os.path.exists(pptx_path):
        print(f"❌ PowerPoint file not found: {pptx_path}")
        return False

    try:
        prs = Presentation(pptx_path)
        print(f"✅ PowerPoint file loaded successfully")
        print(f"   📊 Number of slides: {len(prs.slides)}")

        # Find all merge fields
        merge_fields = set()
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_fields = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    # Find merge fields in format {{field_name}}
                    fields = re.findall(r'\{\{([^}]+)\}\}', text)
                    slide_fields.extend(fields)
                    merge_fields.update(fields)

            if slide_fields:
                print(f"   📝 Slide {slide_num} merge fields: {slide_fields}")

        print(f"   🏷️  Total unique merge fields: {sorted(merge_fields)}")
        return True

    except Exception as e:
        print(f"❌ Error validating PowerPoint file: {e}")
        return False


def validate_config_file(config_path="tests/fixtures/test_config.json"):
    """Validate the test configuration file"""
    print(f"\n🔍 Validating config file: {config_path}")

    if not os.path.exists(config_path):
        print(f"❌ Config file not found: {config_path}")
        return False

    try:
        with open(config_path, 'r') as f:
            config = json.load(f)

        print(f"✅ Config file loaded successfully")

        # Check required sections
        if 'sheet_configs' in config:
            sheets = list(config['sheet_configs'].keys())
            print(f"   📋 Configured sheets: {sheets}")

        if 'global_settings' in config:
            settings = config['global_settings']
            print(f"   ⚙️  Global settings: {list(settings.keys())}")

        return True

    except Exception as e:
        print(f"❌ Error validating config file: {e}")
        return False


def test_data_extraction():
    """Test if the Excel data can be extracted with current config"""
    print(f"\n🧪 Testing data extraction...")

    try:
        # This assumes your project structure - adjust paths as needed
        import sys
        sys.path.append('src')

        from excel_processor import ExcelProcessor
        from config_manager import ConfigManager

        config_manager = ConfigManager()
        config_manager.load_config("tests/fixtures/test_config.json")

        processor = ExcelProcessor(config_manager)
        data = processor.extract_tables_from_excel("tests/fixtures/sample_excel.xlsx")

        print(f"✅ Data extraction successful!")
        print(f"   📊 Extracted tables: {list(data.keys())}")

        for table_name, table_data in data.items():
            print(f"   📋 {table_name}: {len(table_data)} rows")
            if table_data and len(table_data) > 0:
                print(f"      Sample keys: {list(table_data[0].keys())}")

        return True

    except ImportError as e:
        print(f"⚠️  Cannot test extraction - module import failed: {e}")
        print(f"   (This is normal if you haven't set up the project yet)")
        return False
    except Exception as e:
        print(f"❌ Data extraction test failed: {e}")
        return False


def main():
    """Run all validation tests"""
    print("🔍 Validating Test Fixtures")
    print("=" * 50)

    results = []

    # Validate each component
    results.append(validate_excel_file())
    results.append(validate_powerpoint_file())
    results.append(validate_config_file())
    results.append(test_data_extraction())

    print("\n" + "=" * 50)
    if all(results[:3]):  # First 3 are required
        print("✅ All fixture files validated successfully!")
        print("\n🚀 Ready for testing! You can now:")
        print("   1. Run pytest tests")
        print("   2. Start the local server")
        print("   3. Test with the sample files")
    else:
        print("❌ Some validations failed. Please check the errors above.")

    if not results[3]:
        print("\n💡 Note: Data extraction test skipped - run this after setting up the project")


if __name__ == "__main__":
    main()