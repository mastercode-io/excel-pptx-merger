#!/usr/bin/env python3
"""
Script to create sample Excel and PowerPoint test fixtures
Run this to generate test files for the Excel to PowerPoint Merger project
"""

import os
import re
import pandas as pd
from openpyxl import Workbook
from openpyxl.drawing.image import Image as OpenpyxlImage
from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PILImage
from PIL import ImageDraw, ImageFont
import io


def normalize_key(key):
    """Normalize column key by converting to lowercase and replacing spaces with underscores."""
    if not key:
        return ""
    # Convert to lowercase and replace spaces with underscores
    normalized = key.lower().replace(' ', '_')
    # Remove special characters
    normalized = re.sub(r'[^\w_]', '', normalized)
    return normalized


def normalize_image_filename(sheet_name, idx, format_name):
    """Normalize image filename for consistency and reliability."""
    # Normalize sheet name: lowercase, replace spaces with underscores, remove special chars
    normalized_sheet = re.sub(r'[^\w_]', '', sheet_name.lower().replace(' ', '_'))
    # Ensure format is lowercase and without leading dot
    format_lower = format_name.lower().lstrip('.')
    # Create normalized filename
    return f"{normalized_sheet}_image_{idx}.{format_lower}"


def create_sample_excel(output_path="tests/fixtures/sample_excel.xlsx"):
    """Create a sample Excel file matching the expected structure"""

    # Create workbook and select active worksheet
    wb = Workbook()
    ws = wb.active
    ws.title = "Order Form"

    # Table 1: Client Information (horizontal key-value pairs)
    ws['A1'] = "Client"
    ws['B1'] = "Word Or Image"
    ws['C1'] = "G&S Classes"
    ws['D1'] = "SIC"
    ws['E1'] = "Nature of business"
    ws['F1'] = "Designated Countries"

    ws['A2'] = "Acme Corporation Ltd"
    ws['B2'] = "Word"
    ws['C2'] = "9, 35, 42"
    ws['D2'] = "62012"
    ws['E2'] = "Software Development & Consulting"
    ws['F2'] = "United Kingdom, United States"

    # Table 2: Word Search (starting at row 4)
    ws['A4'] = "Word"
    ws['B4'] = "Search Criteria"
    ws['C4'] = "Remarks"

    ws['A5'] = "ACME"
    ws['B5'] = "Identical/Similar"
    ws['C5'] = "Company name search"

    ws['A6'] = "SOFTWARE"
    ws['B6'] = "Phonetic"
    ws['C6'] = "Industry term"

    ws['A7'] = "SOLUTIONS"
    ws['B7'] = "Similar"
    ws['C7'] = "Service description"

    # Table 3: Image Search (starting at row 9)
    ws['A9'] = "Image"
    ws['B9'] = "Search Criteria"
    ws['C9'] = "Image Class.Division.Subdivision"

    ws['A10'] = "[Image will be here]"
    ws['B10'] = "Logo/Brand search"
    ws['C10'] = "26.01.01"

    ws['A11'] = "[Image will be here]"
    ws['B11'] = "Symbol search"
    ws['C11'] = "26.01.03"

    # Create sample images and embed them
    create_and_embed_images(wb, ws)

    # Save the Excel file
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    wb.save(output_path)
    print(f"✅ Created sample Excel file: {output_path}")


def create_and_embed_images(workbook, worksheet):
    """Create sample images and embed them in the Excel file"""

    # Create temporary directory for images
    temp_dir = "temp_images"
    os.makedirs(temp_dir, exist_ok=True)

    try:
        # Create first sample image (logo-style)
        img1 = PILImage.new('RGB', (150, 80), color='#4472C4')
        draw1 = ImageDraw.Draw(img1)

        # Try to use a font, fall back to default if not available
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()

        draw1.text((75, 40), "LOGO", fill='white', anchor='mm', font=font)
        # Use normalized filename
        img1_filename = normalize_image_filename("Order Form", 1, "png")
        img1_path = os.path.join(temp_dir, img1_filename)
        img1.save(img1_path)

        # Create second sample image (symbol-style)
        img2 = PILImage.new('RGB', (120, 120), color='#70AD47')
        draw2 = ImageDraw.Draw(img2)
        # Draw a simple geometric shape
        draw2.ellipse([20, 20, 100, 100], fill='white', outline='#70AD47', width=3)
        draw2.text((60, 60), "®", fill='#70AD47', anchor='mm', font=font)
        # Use normalized filename
        img2_filename = normalize_image_filename("Order Form", 2, "png")
        img2_path = os.path.join(temp_dir, img2_filename)
        img2.save(img2_path)

        # Embed images in Excel
        try:
            # Add first image at row 10 (Image search table)
            img1_xl = OpenpyxlImage(img1_path)
            img1_xl.anchor = 'A10'
            img1_xl.width = 100
            img1_xl.height = 60
            worksheet.add_image(img1_xl)

            # Add second image at row 11
            img2_xl = OpenpyxlImage(img2_path)
            img2_xl.anchor = 'A11'
            img2_xl.width = 80
            img2_xl.height = 80
            worksheet.add_image(img2_xl)

            print("✅ Embedded sample images in Excel file")
            print(f"   - Image 1: {img1_filename} at cell A10")
            print(f"   - Image 2: {img2_filename} at cell A11")

        except Exception as e:
            print(f"⚠️  Warning: Could not embed images in Excel: {e}")

    except Exception as e:
        print(f"⚠️  Warning: Could not create sample images: {e}")


def cleanup_temp_files():
    """Clean up temporary files created during the process"""
    try:
        import shutil
        temp_dir = "temp_images"
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
            print(f"✅ Cleaned up temporary files in {temp_dir}")
    except Exception as e:
        print(f"⚠️  Warning: Could not clean up temporary files: {e}")


def create_sample_powerpoint(output_path="tests/fixtures/sample_template.pptx"):
    """Create a sample PowerPoint template with merge fields"""

    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Trademark Search Report"
    subtitle.text = "Client: {{order_form.client_info.client_name}}"

    # Slide 2: Client Information
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]

    title2.text = "Client Information"
    # Updated field names to match normalized keys
    content2.text = """Client Name: {{order_form.client_info.client_name}}
Search Type: {{order_form.client_info.search_type}}
G&S Classes: {{order_form.client_info.gs_classes}}
SIC Code: {{order_form.client_info.sic_code}}
Business Nature: {{order_form.client_info.business_nature}}
Target Countries: {{order_form.client_info.countries}}"""

    # Slide 3: Word Search Results
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]

    title3.text = "Word Search Results"
    content3.text = """Search performed for the following terms:

Word: {{order_form.word_search.0.word}}
Criteria: {{order_form.word_search.0.search_criteria}}
Remarks: {{order_form.word_search.0.remarks}}

Word: {{order_form.word_search.1.word}}
Criteria: {{order_form.word_search.1.search_criteria}}
Remarks: {{order_form.word_search.1.remarks}}

Word: {{order_form.word_search.2.word}}
Criteria: {{order_form.word_search.2.search_criteria}}
Remarks: {{order_form.word_search.2.remarks}}"""

    # Slide 4: Image Search Results
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title textbox
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Image Search Results"

    # Add content textbox
    content_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2))
    content_frame = content_box.text_frame
    # Updated field names to match normalized keys
    content_frame.text = """Image Search Criteria:

{{order_form.image_search.0.search_criteria}}
Class: {{order_form.image_search.0.image_classification}}

{{order_form.image_search.1.search_criteria}}
Class: {{order_form.image_search.1.image_classification}}"""

    # Add image placeholders using the new format
    img1_placeholder = slide4.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(3), Inches(2))
    img1_frame = img1_placeholder.text_frame
    # Use the image placeholder pattern from config
    img1_frame.text = "{{image:order_form_image_1.png}}"

    img2_placeholder = slide4.shapes.add_textbox(Inches(5.5), Inches(4), Inches(3), Inches(2))
    img2_frame = img2_placeholder.text_frame
    # Use the image placeholder pattern from config
    img2_frame.text = "{{image:order_form_image_2.png}}"

    # Slide 5: Summary
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]

    title5.text = "Search Summary"
    # Updated field names to match normalized keys
    content5.text = """Search completed for {{order_form.client_info.client_name}}

Business Type: {{order_form.client_info.business_nature}}
SIC Code: {{order_form.client_info.sic_code}}
Target Markets: {{order_form.client_info.countries}}

This report covers both word and image searches
as requested for trademark clearance purposes."""

    # Save PowerPoint file
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ Created sample PowerPoint template: {output_path}")


def create_test_files():
    """Create all test fixture files"""
    print("Creating test fixture files...")
    print("=" * 50)

    # Create sample Excel file
    create_sample_excel()

    # Create sample PowerPoint template
    create_sample_powerpoint()
    
    # Clean up temporary files after all files are created
    cleanup_temp_files()

    print("=" * 50)
    print("✅ All test fixtures created successfully!")
    print("\nNext steps:")
    print("1. Run the tests: uv run pytest tests/ -v")
    print("2. Test the API: uv run python scripts/run_local_server.py")
    print("3. Try the sample files:")
    print("   curl -X POST http://localhost:8080/api/v1/merge \\")
    print("     -F 'excel_file=@tests/fixtures/sample_excel.xlsx' \\")
    print("     -F 'pptx_file=@tests/fixtures/sample_template.pptx' \\")
    print("     -F 'config=@tests/fixtures/default_config.json' \\")
    print("     -o 'sample_output.pptx'")


if __name__ == "__main__":
    create_test_files()
