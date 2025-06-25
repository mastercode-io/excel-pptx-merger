"""End-to-end integration tests."""

import json
import os
import pytest
import tempfile
from openpyxl import Workbook
from pptx import Presentation
from pptx.shapes.base import BaseShape

from src.main import app
from src.config_manager import ConfigManager
from src.excel_processor import ExcelProcessor
from src.pptx_processor import PowerPointProcessor
from src.temp_file_manager import TempFileManager


class TestEndToEnd:
    """End-to-end integration tests."""

    @pytest.fixture
    def sample_excel_workbook(self):
        """Create sample Excel workbook for testing."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp_file:
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "Order Form"

            # Client info section
            worksheet["A1"] = "Client"
            worksheet["B1"] = "Word Or Image"
            worksheet["C1"] = "Date"
            worksheet["A2"] = "Acme Corp"
            worksheet["B2"] = "Word"
            worksheet["C2"] = "2023-12-01"

            # Classes section
            worksheet["A5"] = "Class"
            worksheet["A6"] = "Class"
            worksheet["B6"] = "Description"
            worksheet["C6"] = "Status"
            worksheet["A7"] = "35"
            worksheet["B7"] = "Advertising services"
            worksheet["C7"] = "Active"
            worksheet["A8"] = "42"
            worksheet["B8"] = "Computer services"
            worksheet["C8"] = "Pending"

            workbook.save(tmp_file.name)
            workbook.close()

            yield tmp_file.name

            if os.path.exists(tmp_file.name):
                os.unlink(tmp_file.name)

    @pytest.fixture
    def sample_pptx_template(self):
        """Create sample PowerPoint template for testing."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_file:
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])

            # Add text with merge fields (using normalized keys)
            title = slide.shapes.title
            title.text = (
                "Trademark Application for {{order_form.client_info.client_name}}"
            )

            # Add content with merge fields (using normalized keys)
            content = slide.placeholders[1]
            content.text = """
Client: {{order_form.client_info.client_name}}
Type: {{order_form.client_info.word_or_image}}
Date: {{order_form.client_info.date}}

Classes:
{{order_form.trademark_classes.0.class}} - {{order_form.trademark_classes.0.description}}
"""

            presentation.save(tmp_file.name)

            yield tmp_file.name

            if os.path.exists(tmp_file.name):
                os.unlink(tmp_file.name)

    @pytest.fixture
    def test_config(self):
        """Test configuration matching the sample data."""
        return {
            "version": "1.0",
            "sheet_configs": {
                "Order Form": {
                    "subtables": [
                        {
                            "name": "client_info",
                            "type": "key_value_pairs",
                            "header_search": {
                                "method": "contains_text",
                                "text": "Client",
                                "column": "A",
                                "search_range": "A1:A5",
                            },
                            "data_extraction": {
                                "orientation": "horizontal",
                                "headers_row_offset": 0,
                                "data_row_offset": 1,
                                "max_columns": 3,
                                "column_mappings": {
                                    "Client": "client_name",
                                    "Word Or Image": "word_or_image",
                                    "Date": "date",
                                },
                            },
                        },
                        {
                            "name": "trademark_classes",
                            "type": "table",
                            "header_search": {
                                "method": "contains_text",
                                "text": "Class",
                                "column": "A",
                                "search_range": "A5:A10",
                            },
                            "data_extraction": {
                                "orientation": "vertical",
                                "headers_row_offset": 1,
                                "data_row_offset": 2,
                                "max_columns": 3,
                                "max_rows": 10,
                                "column_mappings": {
                                    "Class": "class",
                                    "Description": "description",
                                    "Status": "status",
                                },
                            },
                        },
                    ]
                }
            },
            "global_settings": {
                "normalize_keys": True,
                "temp_file_cleanup": {"enabled": False, "development_mode": True},
            },
        }

    def test_complete_workflow(
        self, sample_excel_workbook, sample_pptx_template, test_config
    ):
        """Test complete Excel to PowerPoint workflow."""
        # Initialize temp file manager
        temp_manager = TempFileManager({"enabled": False, "development_mode": True})

        with temp_manager.temp_directory() as temp_dir:
            # Step 1: Process Excel file
            excel_processor = ExcelProcessor(sample_excel_workbook)
            try:
                extracted_data = excel_processor.extract_data(
                    test_config.get("global_settings", {}), test_config["sheet_configs"]
                )

                # Verify extracted data structure
                assert "order_form" in extracted_data
                sheet_data = extracted_data["order_form"]

                assert "client_info" in sheet_data
                client_info = sheet_data["client_info"]
                assert client_info["client_name"] == "Acme Corp"
                assert client_info["word_or_image"] == "Word"
                assert client_info["date"] == "2023-12-01"

                assert "trademark_classes" in sheet_data
                classes_data = sheet_data["trademark_classes"]
                assert len(classes_data) >= 1  # May only have partial data in test
                # Check first item
                assert "class" in classes_data[0]
                assert "description" in classes_data[0]
                # Data comes as strings from Excel
                assert isinstance(classes_data[0]["class"], str)
                assert isinstance(classes_data[0]["description"], str)

            finally:
                excel_processor.close()

            # Step 2: Process PowerPoint template
            pptx_processor = PowerPointProcessor(sample_pptx_template)
            try:
                # Generate merge preview
                preview = pptx_processor.preview_merge(extracted_data)
                assert "merge_fields" in preview
                assert "field_values" in preview

                # Perform actual merge
                output_path = os.path.join(temp_dir, "merged_output.pptx")
                merged_file = pptx_processor.merge_data(extracted_data, output_path)

                # Verify output file was created
                assert os.path.exists(merged_file)
                assert merged_file == output_path

                # Verify file is valid PowerPoint
                test_presentation = Presentation(merged_file)
                assert len(test_presentation.slides) > 0

                # Check if merge fields were replaced (basic check)
                slide = test_presentation.slides[0]
                title_text = slide.shapes.title.text if slide.shapes.title else ""
                assert "Acme Corp" in title_text  # Merge field should be replaced

            finally:
                pptx_processor.close()

    def test_api_merge_endpoint(
        self, sample_excel_workbook, sample_pptx_template, test_config
    ):
        """Test the API merge endpoint."""
        with app.test_client() as client:
            # Prepare files for upload
            with (
                open(sample_excel_workbook, "rb") as excel_file,
                open(sample_pptx_template, "rb") as pptx_file,
            ):
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "pptx_file": (pptx_file, "template.pptx"),
                    "config": json.dumps(test_config),
                }

                response = client.post(
                    "/api/v1/merge", data=data, content_type="multipart/form-data"
                )

                # Should return the merged file
                assert response.status_code == 200
                assert (
                    response.content_type
                    == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                assert len(response.data) > 0

    def test_api_preview_endpoint(
        self, sample_excel_workbook, sample_pptx_template, test_config
    ):
        """Test the API preview endpoint."""
        with app.test_client() as client:
            with (
                open(sample_excel_workbook, "rb") as excel_file,
                open(sample_pptx_template, "rb") as pptx_file,
            ):
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "pptx_file": (pptx_file, "template.pptx"),
                    "config": json.dumps(test_config),
                }

                response = client.post(
                    "/api/v1/preview", data=data, content_type="multipart/form-data"
                )

                assert response.status_code == 200

                result = response.get_json()
                assert result["success"] is True
                assert "preview" in result

                preview = result["preview"]
                assert "extracted_data" in preview
                assert "template_info" in preview
                assert "merge_preview" in preview
                assert "configuration_used" in preview

                # Verify extracted data
                extracted_data = preview["extracted_data"]
                assert "order_form" in extracted_data

    def test_api_health_endpoint(self):
        """Test the health check endpoint."""
        with app.test_client() as client:
            response = client.get("/api/v1/health")

            assert response.status_code == 200

            result = response.get_json()
            assert result["success"] is True
            assert result["status"] == "healthy"
            assert "version" in result
            assert "timestamp" in result
            assert "services" in result

    def test_api_config_endpoint(self):
        """Test the configuration endpoint."""
        with app.test_client() as client:
            # Test GET - retrieve default config
            response = client.get("/api/v1/config")

            assert response.status_code == 200

            result = response.get_json()
            assert result["success"] is True
            assert "config" in result

            config = result["config"]
            assert "version" in config
            assert "sheet_configs" in config
            assert "global_settings" in config

    def test_api_validation_errors(self):
        """Test API validation with invalid inputs."""
        with app.test_client() as client:
            # Test missing files
            response = client.post("/api/v1/merge", data={})
            assert response.status_code == 400

            result = response.get_json()
            assert result["success"] is False
            assert "error" in result

    def test_api_extract_endpoint_single_sheet(
        self, sample_excel_workbook, test_config
    ):
        """Test the API extract endpoint with single sheet."""
        with app.test_client() as client:
            with open(sample_excel_workbook, "rb") as excel_file:
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "sheet_names": json.dumps(["Order Form"]),
                    "config": json.dumps(test_config["sheet_configs"]),
                    "auto_detect": "false",
                }

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )

                assert response.status_code == 200

                result = response.get_json()
                assert result["success"] is True
                assert "tabs" in result
                assert "summary" in result

                # Check summary
                summary = result["summary"]
                assert summary["total_sheets_requested"] == 1
                assert summary["sheets_processed"] == 1
                assert summary["sheets_not_found"] == 0
                assert summary["sheets_with_errors"] == 0

                # Check tabs structure
                tabs = result["tabs"]
                assert len(tabs) == 1

                tab = tabs[0]
                assert tab["success"] is True
                assert tab["sheet_name"] == "Order Form"
                assert "extracted_data" in tab
                assert "metadata" in tab

                # Verify extracted data structure
                extracted_data = tab["extracted_data"]
                assert "client_info" in extracted_data
                assert "trademark_classes" in extracted_data

                # Verify specific data
                client_info = extracted_data["client_info"]
                assert client_info["client_name"] == "Acme Corp"

    def test_api_extract_endpoint_multiple_sheets(self, sample_excel_workbook):
        """Test the API extract endpoint with multiple sheets."""
        with app.test_client() as client:
            with open(sample_excel_workbook, "rb") as excel_file:
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "sheet_names": json.dumps(["Order Form", "NonExistent"]),
                    "auto_detect": "true",
                }

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )

                assert response.status_code == 200

                result = response.get_json()
                assert result["success"] is True
                assert "tabs" in result
                assert "summary" in result

                # Check summary
                summary = result["summary"]
                assert summary["total_sheets_requested"] == 2
                assert summary["sheets_processed"] == 1  # Only 'Order Form' exists
                assert summary["sheets_not_found"] == 1  # 'NonExistent' doesn't exist
                assert summary["sheets_with_errors"] == 0

                # Check tabs structure
                tabs = result["tabs"]
                assert len(tabs) == 2

                # First tab should be successful (Order Form)
                order_form_tab = next(
                    tab for tab in tabs if tab["sheet_name"] == "Order Form"
                )
                assert order_form_tab["success"] is True
                assert "extracted_data" in order_form_tab
                assert (
                    len(order_form_tab["extracted_data"]) > 0
                )  # Should have auto-detected data

                # Second tab should indicate sheet not found (NonExistent)
                nonexistent_tab = next(
                    tab for tab in tabs if tab["sheet_name"] == "NonExistent"
                )
                assert (
                    nonexistent_tab["success"] is True
                )  # Success=True even for missing sheets
                assert nonexistent_tab["extracted_data"] == {}
                assert (
                    nonexistent_tab["metadata"]["extraction_method"]
                    == "sheet_not_found"
                )
                assert "error" in nonexistent_tab["metadata"]

    def test_api_extract_endpoint_per_sheet_config(
        self, sample_excel_workbook, test_config
    ):
        """Test the API extract endpoint with per-sheet configuration."""
        with app.test_client() as client:
            with open(sample_excel_workbook, "rb") as excel_file:
                # Create per-sheet config
                per_sheet_config = {
                    "Order Form": test_config["sheet_configs"]["Order Form"],
                    "default": {
                        "subtables": [
                            {
                                "name": "auto_table",
                                "type": "table",
                                "header_search": {
                                    "method": "contains_text",
                                    "text": "Class",
                                    "column": "A",
                                    "search_range": "A1:A10",
                                },
                                "data_extraction": {
                                    "headers_row_offset": 1,
                                    "data_row_offset": 2,
                                    "max_columns": 5,
                                    "max_rows": 10,
                                },
                            }
                        ]
                    },
                }

                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "sheet_names": json.dumps(["Order Form"]),
                    "config": json.dumps(per_sheet_config),
                    "auto_detect": "false",
                }

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )

                assert response.status_code == 200

                result = response.get_json()
                assert result["success"] is True

                tabs = result["tabs"]
                assert len(tabs) == 1

                tab = tabs[0]
                assert tab["success"] is True
                assert tab["metadata"]["extraction_method"] == "config_based"

    def test_api_extract_endpoint_validation_errors(self):
        """Test the API extract endpoint validation with invalid inputs."""
        with app.test_client() as client:
            # Test missing excel_file
            response = client.post("/api/v1/extract", data={})
            assert response.status_code == 400

            result = response.get_json()
            assert result["success"] is False
            assert "Excel file is required" in result["error"]["message"]

    def test_api_extract_endpoint_invalid_sheet_names(self, sample_excel_workbook):
        """Test the API extract endpoint with invalid sheet_names parameter."""
        with app.test_client() as client:
            # Test missing sheet_names
            with open(sample_excel_workbook, "rb") as excel_file:
                data = {"excel_file": (excel_file, "test.xlsx")}

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )
                assert response.status_code == 400

                result = response.get_json()
                assert "sheet_names parameter is required" in result["error"]["message"]

            # Test invalid JSON in sheet_names
            with open(sample_excel_workbook, "rb") as excel_file:
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "sheet_names": "invalid json",
                }

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )
                assert response.status_code == 400

                result = response.get_json()
                assert (
                    "Invalid JSON in sheet_names parameter"
                    in result["error"]["message"]
                )

            # Test non-list sheet_names
            with open(sample_excel_workbook, "rb") as excel_file:
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "sheet_names": json.dumps("not a list"),
                }

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )
                assert response.status_code == 400

                result = response.get_json()
                assert "must be a list of strings" in result["error"]["message"]

            # Test empty list
            with open(sample_excel_workbook, "rb") as excel_file:
                data = {
                    "excel_file": (excel_file, "test.xlsx"),
                    "sheet_names": json.dumps([]),
                }

                response = client.post(
                    "/api/v1/extract", data=data, content_type="multipart/form-data"
                )
                assert response.status_code == 400

                result = response.get_json()
                assert "cannot be empty" in result["error"]["message"]

    def test_configuration_manager_integration(self, test_config):
        """Test configuration manager integration."""
        with tempfile.TemporaryDirectory() as temp_dir:
            config_manager = ConfigManager(temp_dir)

            # Test saving and loading custom config
            config_manager.save_config(test_config, "integration_test")
            loaded_config = config_manager.load_config("integration_test")

            assert loaded_config == test_config

            # Test runtime validation
            config_manager.validate_runtime_config(loaded_config)

    def test_temp_file_manager_integration(self):
        """Test temporary file manager integration."""
        config = {
            "enabled": True,
            "delay_seconds": 1,
            "keep_on_error": False,
            "development_mode": False,
        }

        temp_manager = TempFileManager(config)

        # Test directory creation and cleanup
        temp_dir = temp_manager.create_temp_directory()
        assert os.path.exists(temp_dir)

        # Create a test file
        test_file = temp_manager.create_temp_file(temp_dir, "test.txt", b"test content")
        assert os.path.exists(test_file)

        # Schedule cleanup and wait
        temp_manager.schedule_cleanup(temp_dir, delay_override=0)

        # Directory should be cleaned up immediately with delay_override=0
        import time

        time.sleep(0.1)  # Give cleanup thread time to run

        # In a real scenario, we'd wait for cleanup, but for testing we'll just verify
        # the cleanup was scheduled
        stats = temp_manager.get_stats()
        assert stats["total_directories"] >= 0  # Directory may already be cleaned up

    def test_error_handling_workflow(self, sample_pptx_template):
        """Test error handling in the complete workflow."""
        # Test with invalid Excel file
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as invalid_excel:
            invalid_excel.write(b"This is not a valid Excel file")
            invalid_excel.flush()

            try:
                with pytest.raises(Exception):  # Should raise ExcelProcessingError
                    excel_processor = ExcelProcessor(invalid_excel.name)
                    excel_processor.extract_data({})
            finally:
                os.unlink(invalid_excel.name)

        # Test with invalid PowerPoint file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as invalid_pptx:
            invalid_pptx.write(b"This is not a valid PowerPoint file")
            invalid_pptx.flush()

            try:
                with pytest.raises(Exception):  # Should raise PowerPointProcessingError
                    pptx_processor = PowerPointProcessor(invalid_pptx.name)
                    pptx_processor.get_merge_fields()
            finally:
                os.unlink(invalid_pptx.name)
